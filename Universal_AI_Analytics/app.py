"""
app.py – Universal Analytics + Digital Marketing FastAPI Backend
================================================================

Three-phase pipeline:
  Phase 1 → LangGraph strategic planning (5 nodes, Gemini)
  HITL    → Human approval checkpoint
  Phase 2 → Analytics CrewAI (5 agents: Engineer, Scientist, LabelEncoder, Analyst, CDO)
             + Gemini enrichment of analytics report
  Phase 3 → Digital Marketing CrewAI (4 agents) fed analytics context

User provides a Database URI (any SQLAlchemy-compatible DB) — the system
auto-discovers the schema and adapts all analysis, ML models, and marketing
campaigns to whatever data is in the connected database.

7-tab Bootstrap 5 dark UI:
  Console | Analytics Report | Charts | Marketing Report | Content | Downloads | Conversation

Run:
    python main.py   →  http://localhost:8002
"""

import os
import sys

os.environ["OTEL_SDK_DISABLED"]         = "true"
os.environ["CREWAI_TELEMETRY_OPT_OUT"] = "true"

import re
import json
import uuid
import queue
import asyncio
import threading
from datetime import datetime
from pathlib import Path
from typing import List, Optional
import pandas as pd

BASE_DIR = Path(__file__).parent
sys.path.insert(0, str(BASE_DIR))

for _d in ("outputs/charts", "outputs/reports", "outputs/models",
           "outputs/videos", "outputs/content", "outputs/audit", "outputs/uploads", "static"):
    (BASE_DIR / _d).mkdir(parents=True, exist_ok=True)

AUDIT_DIR = BASE_DIR / "outputs" / "audit"

from fastapi import FastAPI, HTTPException, File, UploadFile
from fastapi.responses import StreamingResponse, FileResponse
from fastapi.staticfiles import StaticFiles
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel

app = FastAPI(
    title="Universal Analytics + Digital Marketing Multi-Agent System",
    description="Connect any database — AI agents auto-discover, analyse, and build marketing campaigns",
    version="2.0.0",
)

app.add_middleware(CORSMiddleware, allow_origins=["*"], allow_credentials=True,
                   allow_methods=["*"], allow_headers=["*"])

app.mount("/static",  StaticFiles(directory=str(BASE_DIR / "static")),  name="static")
app.mount("/outputs", StaticFiles(directory=str(BASE_DIR / "outputs")), name="outputs")

# ── In-memory store ───────────────────────────────────────────────────────────
tasks: dict       = {}
log_queues: dict  = {}
hitl_events: dict = {}
_analysis_lock    = threading.Lock()
_uploaded_files: dict = {}

# ── Security patterns ────────────────────────────────────────────────────────
_PII_PATTERNS = [
    (re.compile(r'\b[A-Za-z0-9._%+\-]+@[A-Za-z0-9.\-]+\.[A-Za-z]{2,}\b'), '[EMAIL]'),
    (re.compile(r'(\+62|62|0)[\s.\-]?8[1-9][0-9]{6,10}\b'), '[PHONE]'),
    (re.compile(r'\b\d{16}\b'), '[NIK/ACCT]'),
    (re.compile(r'\b(?:\d{4}[\s\-]?){3}\d{4}\b'), '[CARD]'),
    (re.compile(r'\b\d{10,15}\b'), '[ACCT]'),
]

_SQL_RE = re.compile(
    r'(union\s+(all\s+)?select)|(drop\s+(table|database))'
    r'|(insert\s+into\s+\w)|(delete\s+from\s+\w)|(update\s+\w+\s+set\s)'
    r'|(exec\s*\(|xp_cmdshell)|(\bor\b\s+[\'"]?\d[\'"]?\s*=\s*[\'"]?\d)'
    r'|(--\s*[\r\n]|;\s*--)|(\/\*[\s\S]*?\*\/)',
    re.IGNORECASE,
)

_GUARDRAIL_RE = re.compile(
    r'\b(transfer|send|wire)\s+(money|fund|cash|rp|usd|idr)'
    r'|\b(withdraw|deposit)\s+(from|to|my)'
    r'|\bmy\s+(account\s+(number|balance|pin|password)|card\s+number)'
    r'|\b(ignore|forget|disregard)\s+(previous|all|above|your)\s+(instruction|rule|prompt)'
    r'|\b(jailbreak|prompt.injection|bypass.filter|bypass.security)'
    r'|\b(act\s+as\s+a\s+|you\s+are\s+now\s+a|pretend\s+(you\s+are|to\s+be))'
    r'|\b(override|reset)\s+your\s+(system|instruction|role)',
    re.IGNORECASE,
)


def _redact_pii(text: str):
    found = False
    for pattern, placeholder in _PII_PATTERNS:
        new_text = pattern.sub(placeholder, text)
        if new_text != text:
            found = True
            text = new_text
    return text, found


def _is_sql_injection(text: str) -> bool:
    return bool(_SQL_RE.search(text))


def _is_guardrail_violation(text: str) -> bool:
    return bool(_GUARDRAIL_RE.search(text))


def _audit_log(entry: dict):
    try:
        with open(AUDIT_DIR / "conversation_audit.jsonl", "a", encoding="utf-8") as fh:
            fh.write(json.dumps(entry, ensure_ascii=False) + "\n")
    except Exception:
        pass


# ── Pydantic models ───────────────────────────────────────────────────────────

class CombinedRequest(BaseModel):
    # Data sources — all optional; system auto-detects what's available
    database_uri:      str       = ""   # SQLAlchemy-compatible URI (optional)
    uploaded_file_ids: List[str] = []   # IDs returned by /api/upload (optional)
    # Analytics inputs
    analysis_request:  str       = ""
    use_langgraph:     bool      = True
    # Marketing inputs
    brand_name:        str = "Brand"
    industry:          str = "General"
    target_audience:   str = "General audience"
    campaign_goals:    str = "Grow retention and engagement"
    budget:            str = "Not specified"
    competitors:       str = "Not specified"
    campaign_type:     str = "Retention & Growth"


class ChatRequest(BaseModel):
    message: str


class ApproveRequest(BaseModel):
    abort: bool = False


# ── Stdout capture ────────────────────────────────────────────────────────────

class _TeeWriter:
    """Capture stdout and forward to the SSE log queue, suppressing CrewAI internal noise."""

    # Substrings that mark noisy CrewAI/ReAct internals — skip any write containing these
    _NOISE = (
        "You ONLY have access to the following tools",
        '"additionalProperties"',
        '"title": "Table Name"',
        '"title": "Target Column"',
        "Tool Name:",
        "Tool Arguments:",
        "Tool Description:",
        '"type": "object"',
        '"type": "string"',
        '"required":',
        '"properties":',
        "IMPORTANT: Use the following format",
        "Use the following format in your response",
        "Thought: you should always think",
        "Action Input: the input to the action",
        "Observation: the result of the action",
        "Once all necessary information is gathered",
        "I now know the final answer",
    )

    def __init__(self, task_id: str, original):
        self._tid          = task_id
        self._original     = original
        self._skip_block   = False   # True while inside a tool-schema dump block

    def write(self, text: str):
        self._original.write(text)
        self._original.flush()

        # Enter skip-block when the tool-list header appears
        if "You ONLY have access to the following tools" in text:
            self._skip_block = True

        # Exit skip-block when a new meaningful section starts
        if self._skip_block:
            clean = text.strip()
            if clean.startswith(("[", "━", "✅", "🔧", "Phase", "Error", "WARNING")) or \
               any(clean.startswith(p) for p in ("Thought:", "Final Answer:", "Action:", "[Content", "[Bank", "[Crew")):
                self._skip_block = False
            else:
                return  # drop noisy schema/block content

        stripped = text.rstrip()
        if not stripped:
            return

        # Drop individual lines containing known noise patterns
        if any(p in text for p in self._NOISE):
            return

        # Drop lines that are purely raw JSON fragments (common in verbose tool dumps)
        first = stripped.lstrip()
        if first.startswith(('"', '{', '}', '[', ']', '|')) and len(stripped) > 40:
            return

        entry = {"type": "log", "message": stripped, "timestamp": datetime.now().isoformat()}
        tasks[self._tid]["logs"].append(entry)
        try:
            log_queues[self._tid].put_nowait(entry)
        except Exception:
            pass

    def flush(self):
        self._original.flush()

    def isatty(self):
        return False


def _push(task_id: str, msg: str, kind: str = "log"):
    entry = {"type": kind, "message": msg, "timestamp": datetime.now().isoformat()}
    tasks[task_id]["logs"].append(entry)
    try:
        log_queues[task_id].put_nowait(entry)
    except Exception:
        pass


def _strip_md_fences(text: str) -> str:
    text = re.sub(r'^```(?:markdown|md)?\s*\n', '', text.strip())
    text = re.sub(r'\n```\s*$', '', text.strip())
    text = re.sub(r'\n```(?:markdown|md)?\n[\s\S]*?```(?:\n|$)', '\n', text)
    return text.strip()


def _default_analytics_request() -> str:
    return (
        "Perform a comprehensive end-to-end data analysis on the connected data source.\n\n"
        "━━━ STEP 1 — DATA ENGINEERING (Data Engineer Agent) ━━━\n"
        "• Call list_database_tables to discover ALL available tables.\n"
        "• Call profile_database_table on the primary/most relevant table.\n"
        "• Record: row count, column names, data types, null percentages, numeric statistics,\n"
        "  and the value distribution of any binary or categorical target column.\n"
        "• Remove ID/UUID columns that add no predictive value.\n"
        "• Normalize data types (convert text-stored numbers to numeric).\n"
        "• Report data quality findings (missing values, outliers, type mismatches).\n\n"
        "━━━ STEP 2 — MACHINE LEARNING (Data Scientist Agent) ━━━\n"
        "• Inspect the schema and automatically choose the right ML approach:\n"
        "  - Binary target column found → train classification model (churn, fraud, or credit risk)\n"
        "  - Continuous target column found → regression analysis\n"
        "  - Date/timestamp column found → run time_series_forecast in addition to other models\n"
        "  - No clear target → run customer_segmentation only\n"
        "• Always run customer_segmentation (4 segments) regardless of other models.\n"
        "• Report AUC-ROC (or RMSE for regression), accuracy, and the top 5 most important\n"
        "  features with plain-English business explanations of why each matters.\n"
        "• Translate ML outputs into business language — avoid jargon.\n\n"
        "━━━ STEP 3 — DATA VISUALIZATION (Data Analyst Agent) ━━━\n"
        "• Generate at least 5 individual charts using ACTUAL column names from the data:\n"
        "  1. Target/outcome distribution — pie chart\n"
        "  2. Key numeric feature — histogram\n"
        "  3. Numeric feature grouped by target — histplot\n"
        "  4. Two-numeric relationship — scatter plot\n"
        "  5. Category vs. numeric — bar chart\n"
        "  6. Full feature correlation — heatmap (use the _encoded table)\n"
        "• After each chart, call Gemini AI vision to analyse what the chart reveals.\n"
        "• Extract 2–3 actionable business insights per chart.\n"
        "• Do NOT generate a dashboard — individual charts only.\n\n"
        "━━━ STEP 4 — REPORTING (CDO/Manager Agent) ━━━\n"
        "• Generate PDF and PowerPoint executive reports.\n"
        "• Write an Executive Summary with the top 5 findings, each with:\n"
        "  - A specific metric (e.g. '34% of Segment A churned in Q3')\n"
        "  - A clear business implication\n"
        "  - A concrete recommendation with an owner, KPI target, and timeline\n"
        "• End with a numbered action plan: [N]. Action — Owner — Deadline — Success Metric\n"
    )


def _list_files(directory: str, extensions: list) -> list:
    p = BASE_DIR / directory
    if not p.exists():
        return []
    items = []
    for ext in extensions:
        for f in sorted(p.glob(f"*{ext}"), key=os.path.getmtime, reverse=True):
            rel = f.relative_to(BASE_DIR).as_posix()
            items.append({
                "name":     f.name,
                "url":      f"/{rel}",
                "size":     f.stat().st_size,
                "ext":      ext.lstrip("."),
                "modified": datetime.fromtimestamp(f.stat().st_mtime).isoformat(),
            })
    return items


# ── Document text extraction helpers ─────────────────────────────────────────

def _extract_pdf_text(path: str) -> str:
    try:
        import pdfplumber
        parts = []
        with pdfplumber.open(path) as pdf:
            for page in pdf.pages:
                t = page.extract_text()
                if t:
                    parts.append(t)
        return "\n\n".join(parts)
    except Exception:
        pass
    try:
        import pypdf
        r = pypdf.PdfReader(path)
        return "\n\n".join(p.extract_text() or "" for p in r.pages)
    except Exception:
        return ""


def _extract_word_text(path: str) -> str:
    try:
        from docx import Document
        doc = Document(path)
        return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except Exception:
        return ""


def _extract_pptx_text(path: str) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(path)
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            texts = [sh.text.strip() for sh in slide.shapes if hasattr(sh, "text") and sh.text.strip()]
            if texts:
                parts.append(f"[Slide {i}]\n" + "\n".join(texts))
        return "\n\n".join(parts)
    except Exception:
        return ""


def _build_adaptive_request(
    has_db: bool,
    has_documents: bool,
    has_any_data: bool,
    struct_summary: str,
    doc_context: str,
) -> str:
    """Build an analysis prompt that adapts to whatever data sources are available."""
    lines = [
        "Perform a comprehensive analysis using whatever data sources are available below.\n"
        "Adapt your methodology to fit the inputs — not all steps apply in every scenario.\n"
    ]
    if has_db:
        lines += [
            "\n━━━ STEP 1 — DATA ENGINEERING ━━━",
            "• Call list_database_tables to discover all available tables.",
            "• Profile the primary/most relevant table: row count, columns, types, nulls, distributions.",
            "• Remove ID/UUID columns, normalize data types, report data quality issues.",
        ]
        if struct_summary:
            lines.append(f"  Available structured sources:\n{struct_summary}")
        lines += [
            "\n━━━ STEP 2 — MACHINE LEARNING ━━━",
            "• Inspect the schema and choose the appropriate model:",
            "  - Binary target column → classification (churn / fraud / risk)",
            "  - Continuous target → regression",
            "  - Date/time column found → also run time_series_forecast",
            "  - No clear target → customer_segmentation (4 clusters)",
            "• Always run customer_segmentation regardless of other models.",
            "• Report metrics and top feature importances in plain business language.",
            "\n━━━ STEP 3 — VISUALIZATION ━━━",
            "• Generate 5+ individual charts using real column names from the data.",
            "• Analyse each chart with Gemini AI vision — extract 2-3 business insights per chart.",
            "• Do NOT generate a dashboard — individual charts only.",
        ]
    else:
        lines.append("\n━━━ DATA NOTE ━━━\nNo structured database or tabular file was provided. Skip ML and visualization steps.")

    if has_documents:
        lines += [
            "\n━━━ DOCUMENT ANALYSIS ━━━",
            "The following documents have been provided. Extract all key insights, data points,",
            "facts, metrics, and strategic recommendations from them:\n",
            doc_context,
            "\nSummarise the main findings and highlight any quantitative data or actionable insights.",
        ]

    if not has_any_data:
        lines += [
            "\n━━━ RESEARCH MODE — Tavily Web Search ━━━",
            "No data was provided by the user. Use web_search_collect to research:",
            "• Industry benchmarks and market data relevant to this domain",
            "• Competitor landscape and positioning",
            "• Current trends, statistics, and best practices",
            "Build the entire analysis from web research findings.",
        ]
    elif not has_db:
        lines += [
            "\n━━━ SUPPLEMENTAL RESEARCH ━━━",
            "If the document context is insufficient for a complete analysis,",
            "use web_search_collect to fill gaps with industry data and benchmarks.",
        ]

    lines += [
        "\n━━━ STEP 4 — REPORTING ━━━",
        "• Generate PDF and PowerPoint executive reports.",
        "• Executive Summary: top 5 findings with specific metrics.",
        "• Recommendations table: Finding | Recommendation | Owner | Timeline | KPI Target.",
        "• If data was limited, include a section recommending additional data to collect.",
    ]
    return "\n".join(lines)


def _merge_db_sources(sqlite_uri: str, db_uri: str, task_id: str) -> str:
    """Load tables from an external DB into the SQLite file DB and return the SQLite URI."""
    try:
        import pandas as pd
        from sqlalchemy import create_engine, inspect as sa_inspect
        src_engine  = create_engine(db_uri, connect_args={"connect_timeout": 10})
        dest_engine = create_engine(sqlite_uri)
        inspector   = sa_inspect(src_engine)
        tables      = inspector.get_table_names()[:10]
        for tbl in tables:
            try:
                df = pd.read_sql(f"SELECT * FROM {tbl} LIMIT 50000", src_engine)
                dest_name = "db_" + "".join(c if c.isalnum() else "_" for c in tbl)
                df.to_sql(dest_name, dest_engine, if_exists="replace", index=False)
                _push(task_id, f"  Merged DB table '{tbl}' → '{dest_name}' ({len(df)} rows)", "info")
            except Exception as e:
                _push(task_id, f"  Could not merge table '{tbl}': {e}", "warn")
        return sqlite_uri
    except Exception as e:
        _push(task_id, f"  DB merge failed ({e}) — using file data only", "warn")
        return sqlite_uri


# ── Main pipeline ─────────────────────────────────────────────────────────────

def _run_analysis(task_id: str, req: CombinedRequest):
    original_stdout = sys.stdout
    sys.stdout = _TeeWriter(task_id, original_stdout)

    try:
        _push(task_id, "━" * 62, "info")
        _push(task_id, "  ANALYTICS + DIGITAL MARKETING SYSTEM  —  starting", "info")
        _push(task_id, "━" * 62, "info")

        # ── Clear session sidecar files ────────────────────────────────────────
        for _sidecar in (
            BASE_DIR / "outputs" / "charts" / "session_charts.json",
            BASE_DIR / "outputs" / "content" / "session_content.json",
        ):
            try:
                if _sidecar.exists():
                    _sidecar.unlink()
            except Exception:
                pass

        _brand_name = req.brand_name.strip() or "Brand"
        tasks[task_id]["brand_name"] = _brand_name

        # ── Resolve data sources (fully adaptive — all inputs optional) ────
        from config import POSTGRES_URI

        _structured_files = []   # CSV/Excel → SQLite tables
        _document_files   = []   # PDF/Word/PPTX → text context

        for fid in (req.uploaded_file_ids or []):
            info = _uploaded_files.get(fid.strip())
            if not info:
                continue
            if info["file_type"] == "structured":
                _structured_files.append(info)
            elif info["file_type"] == "document":
                _document_files.append(info)

        _has_explicit_db  = bool(req.database_uri.strip())
        _has_structured   = bool(_structured_files)
        _has_documents    = bool(_document_files)
        _has_any_data     = _has_explicit_db or _has_structured or _has_documents

        # ── Build unified SQLite when structured files exist ───────────────
        _active_db: str | None = None
        if _has_structured or _has_explicit_db:
            _db_uri_to_merge = req.database_uri.strip() if _has_explicit_db else None
            if _has_structured:
                # Merge all structured files into one SQLite
                from sqlalchemy import create_engine as _ce_mrg
                ts_mrg     = datetime.now().strftime("%Y%m%d_%H%M%S")
                _sq_path   = BASE_DIR / "outputs" / "uploads" / f"combined_{ts_mrg}.db"
                _dest_eng  = _ce_mrg(f"sqlite:///{_sq_path}")
                for sfi in _structured_files:
                    try:
                        _src_eng = _ce_mrg(sfi["sqlite_uri"])
                        _df_tmp  = pd.read_sql(f"SELECT * FROM {sfi['table_name']}", _src_eng)
                        _df_tmp.to_sql(sfi["table_name"], _dest_eng, if_exists="replace", index=False)
                        _push(task_id, f"  Loaded '{sfi['table_name']}' from {sfi['original_name']} ({len(_df_tmp)} rows)", "info")
                    except Exception as _e:
                        _push(task_id, f"  Skipped '{sfi['original_name']}': {_e}", "warn")
                _active_db = f"sqlite:///{_sq_path}"
                if _db_uri_to_merge:
                    _active_db = _merge_db_sources(_active_db, _db_uri_to_merge, task_id)
            else:
                _active_db = _db_uri_to_merge
        # else: no structured data — document/research-only mode

        if _active_db:
            os.environ["ACTIVE_DB_URI"] = _active_db
        elif "ACTIVE_DB_URI" in os.environ:
            del os.environ["ACTIVE_DB_URI"]

        # ── Build document context string ─────────────────────────────────
        _doc_context = ""
        if _has_documents:
            _parts = []
            for di in _document_files:
                if di.get("extracted_text"):
                    _parts.append(
                        f"### {di['original_name']} ({di['format'].upper()})\n"
                        f"{di['extracted_text']}"
                    )
            _doc_context = "\n\n---\n\n".join(_parts)

        # ── Build structured-file summary for prompt injection ────────────
        _struct_summary = ""
        if _has_structured:
            _struct_summary = "\n".join(
                f"  • {s['original_name']} → table '{s['table_name']}' "
                f"({s['rows']} rows, {len(s['columns'])} cols: {', '.join(s['columns'][:15])})"
                for s in _structured_files
            )
        if _has_explicit_db:
            _struct_summary += f"\n  • Database URI: {req.database_uri.strip()[:80]}…"

        # ── Log data source summary ────────────────────────────────────────
        _source_parts = []
        if _has_explicit_db:   _source_parts.append("Database")
        if _has_structured:    _source_parts.append(f"{len(_structured_files)} structured file(s)")
        if _has_documents:     _source_parts.append(f"{len(_document_files)} document(s)")
        if not _source_parts:  _source_parts.append("No data — Tavily research mode")
        _push(task_id, f"  Data sources: {' + '.join(_source_parts)}", "info")

        # ── Build adaptive analysis request ───────────────────────────────
        if req.analysis_request.strip():
            analysis_request = req.analysis_request.strip()
            if _struct_summary:
                analysis_request += f"\n\n[AVAILABLE STRUCTURED DATA]:\n{_struct_summary}"
            if _doc_context:
                analysis_request += f"\n\n[DOCUMENT CONTEXT (extracted text)]:\n{_doc_context}"
        else:
            analysis_request = _build_adaptive_request(
                has_db         = _has_explicit_db or _has_structured,
                has_documents  = _has_documents,
                has_any_data   = _has_any_data,
                struct_summary = _struct_summary,
                doc_context    = _doc_context,
            )

        # ══════════════════════════════════════════════════════════════════════
        # PHASE 1 — LangGraph Strategic Planning
        # ══════════════════════════════════════════════════════════════════════
        langgraph_plan = None
        if req.use_langgraph:
            print("\n" + "=" * 62)
            print("  PHASE 1 — LangGraph Strategic Planning")
            print("=" * 62)

            from graphs.banking_graph import run_banking_analysis as _lg_run
            graph_state = _lg_run(analysis_request, [])

            langgraph_plan = {
                "analysis_plan":      graph_state.get("analysis_plan", ""),
                "etl_guidance":       graph_state.get("etl_guidance", ""),
                "ml_guidance":        graph_state.get("ml_guidance", ""),
                "analytics_guidance": graph_state.get("analytics_guidance", ""),
                "preliminary_report": graph_state.get("report_content", ""),
            }
            tasks[task_id]["langgraph_plan"] = langgraph_plan
            print("\n  LangGraph strategic planning complete.\n")

            # ── HITL Pause ─────────────────────────────────────────────────────
            _push(task_id, "━" * 58, "info")
            _push(task_id, "  HUMAN IN THE LOOP — Review the strategic plan", "info")
            _push(task_id, "  Approve in the UI to continue to Phase 2.", "info")
            _push(task_id, "  Auto-approving in 10 minutes if no response.", "info")
            _push(task_id, "━" * 58, "info")

            hitl_entry = {
                "type":               "hitl_pause",
                "task_id":            task_id,
                "plan_preview":       langgraph_plan.get("analysis_plan", "")[:800],
                "preliminary_report": langgraph_plan.get("preliminary_report", "")[:2000],
                "timestamp":          datetime.now().isoformat(),
                "auto_approve_secs":  600,
            }
            tasks[task_id]["logs"].append(hitl_entry)
            try:
                log_queues[task_id].put_nowait(hitl_entry)
            except Exception:
                pass

            tasks[task_id]["status"] = "awaiting_approval"
            approved = hitl_events[task_id].wait(timeout=600)

            if tasks[task_id].get("hitl_aborted", False):
                raise RuntimeError("Analysis aborted by user during HITL review.")

            if not approved:
                _push(task_id, "  HITL: 10-minute timeout — auto-approving.", "warn")
            else:
                _push(task_id, "  HITL: Approved — proceeding to Phase 2.", "success")

            tasks[task_id]["status"] = "running"

        # ══════════════════════════════════════════════════════════════════════
        # PHASE 2 — Analytics CrewAI
        # ══════════════════════════════════════════════════════════════════════
        print("\n" + "=" * 62)
        print("  PHASE 2 — Analytics CrewAI Multi-Agent Execution")
        print("=" * 62)

        from crew_banking import run_banking_crew
        banking_result      = run_banking_crew(analysis_request, langgraph_plan)
        banking_crew_output = banking_result.get("crew_output", "")
        analyst_output      = banking_result.get("analyst_output", "")
        scientist_output    = banking_result.get("scientist_output", "")

        # Load banking blog sections (charts with Gemini analysis)
        banking_blog_sections: list = []
        _meta_file = BASE_DIR / "outputs" / "charts" / "session_charts.json"
        if _meta_file.exists():
            try:
                banking_blog_sections = json.loads(_meta_file.read_text(encoding="utf-8"))
                # Deduplicate
                _seen: set = set()
                _deduped: list = []
                for _s in banking_blog_sections:
                    _url = _s.get("chart_url", "")
                    if _url not in _seen:
                        _seen.add(_url)
                        _deduped.append(_s)
                banking_blog_sections = _deduped
                _push(task_id, f"  Banking charts loaded: {len(banking_blog_sections)}", "info")
            except Exception as _e:
                _push(task_id, f"  Banking chart metadata read failed: {_e}", "warn")

        banking_charts  = _list_files("outputs/charts", [".png"])
        banking_reports = _list_files("outputs/reports", [".pdf", ".pptx", ".md"])

        # Build banking markdown report
        banking_report = banking_crew_output if len(banking_crew_output) > 50 else \
            langgraph_plan.get("preliminary_report", "# Banking Analysis Complete") if langgraph_plan else \
            "# Banking Analysis Complete"

        # ── Ollama enrichment of banking report ──────────────────────────────
        banking_enhanced = banking_report
        try:
            from langchain_community.chat_models import ChatOllama
            from langchain_core.messages import HumanMessage

            if banking_report and len(banking_report) > 100:
                chart_analyses = "\n\n".join(
                    f"**{s['title']}**:\n{s['gemini_analysis']}"
                    for s in banking_blog_sections
                    if s.get("gemini_analysis")
                )
                enhance_prompt = (
                    f"You are a senior strategy consultant preparing a C-suite briefing "
                    f"for a {req.industry} organisation.\n\n"
                    "## Draft Analytics Report:\n" + banking_report + "\n\n"
                    + (f"## Chart Evidence (for your reference ONLY — do NOT copy into the report):\n{chart_analyses}\n\n" if chart_analyses else "")
                    + "## Your Task — Targeted Enrichment:\n"
                    "Produce the final CEO-ready version (600–900 words). "
                    "Keep Executive Summary max 4 bullets. "
                    "Keep Key Findings max 5 points — you may reference a specific number from the "
                    "Chart Evidence above but do NOT reproduce the full chart analysis text. "
                    "Enrich Recommendations with Owner|KPI Target|30-day|60-day|90-day table. "
                    "Replace Next Steps with numbered action plan [N]. Action — Owner — Deadline — Metric. "
                    "IMPORTANT: Do NOT include chart descriptions, visualization analyses, or image "
                    "references in the report body — those are rendered separately in the Chart Analysis "
                    "section below the report. "
                    "Output the complete report only — no preamble, no chart analysis blocks."
                )
                _push(task_id, "  Ollama: enriching analytics report…", "info")
                _ollama = ChatOllama(model="qwen3.5:cloud")
                response = _ollama.invoke([HumanMessage(content=enhance_prompt)])
                enhanced = _strip_md_fences(response.content)
                if enhanced and len(enhanced) > 200:
                    banking_enhanced = enhanced
                    _push(task_id, "  Ollama: analytics report enhanced ✓", "success")
        except Exception as _gem_err:
            _push(task_id, f"  Ollama enrichment (non-critical): {_gem_err}", "warn")

        tasks[task_id].update({
            "banking_report":        banking_enhanced,
            "banking_blog_sections": banking_blog_sections,
            "banking_charts":        banking_charts,
            "analyst_output":        analyst_output,
            "scientist_output":      scientist_output,
        })

        _push(task_id, "━" * 62, "success")
        _push(task_id, f"  Phase 2 complete! Charts: {len(banking_charts)}  Reports: {len(banking_reports)}", "success")

        # ══════════════════════════════════════════════════════════════════════
        # PHASE 3 — Marketing LangGraph Planning + Digital Marketing CrewAI
        # ══════════════════════════════════════════════════════════════════════
        print("\n" + "=" * 62)
        print("  PHASE 3 — Marketing LangGraph + CrewAI Execution")
        print("=" * 62)
        _push(task_id, "━" * 62, "info")
        _push(task_id, "  PHASE 3 — Analytics handing findings to Marketing team…", "info")
        _push(task_id, "━" * 62, "info")

        # Analytics context passed to all marketing agents.
        # Priority: Data Analyst raw output first, then CDO executive summary as supplement.
        if analyst_output:
            banking_context = (
                f"[DATA ANALYST FINDINGS]:\n{analyst_output[:2000]}\n\n"
                f"[CDO EXECUTIVE SUMMARY]:\n{banking_enhanced[:1000]}"
            )
        else:
            banking_context = banking_enhanced[:2000] if banking_enhanced else ""
        analytics_context = banking_context  # unified variable name

        # ── Build chart analysis summary (Data Analyst's findings) ───────────
        _chart_analyses = "\n\n".join(
            f"**{s['title']}**:\n{s.get('gemini_analysis', '')}"
            for s in banking_blog_sections
            if s.get("gemini_analysis")
        )[:2500]

        # ── Generate data-driven campaign brief from Data Analyst output ─────
        # The campaign request is NOT user-typed — it comes entirely from the
        # banking analytics: churn model findings, customer segments, and the
        # Data Analyst's chart insights. The marketing team uses this to design
        # specific, named campaigns (cashback, loyalty programs, etc.).
        _push(task_id, "  Translating analyst findings → data-driven campaign brief…", "info")
        campaign_request = ""
        derived_audience = req.target_audience or "Audience segments identified in analytics"
        derived_goals    = req.campaign_goals   or "Grow retention and engagement"

        try:
            from langchain_community.chat_models import ChatOllama
            from langchain_core.messages import HumanMessage as _HMsg

            if banking_enhanced:
                _analyst_section   = f"## Data Analyst Report:\n{analyst_output[:2000]}\n\n"       if analyst_output   else ""
                _scientist_section = f"## Data Scientist ML Results:\n{scientist_output[:1000]}\n\n" if scientist_output else ""
                _chart_section     = f"## Data Analyst Chart Insights:\n{_chart_analyses}\n\n"      if _chart_analyses  else ""

                _handoff_prompt = (
                    f"You are a Chief Data Officer handing off data analytics results "
                    f"to the Digital Marketing team.\n\n"
                    f"Brand: {_brand_name} | Industry: {req.industry} | Budget: {req.budget}\n\n"
                    "Write a detailed **Digital Marketing Campaign Brief** based strictly on "
                    "the Data Analyst's findings below. The brief must include:\n\n"
                    "1. **At-Risk or High-Opportunity Segments** — name each segment with specific "
                    "numbers from the data (e.g. '34% of Segment A churned', "
                    "'Segment B has 3× higher lifetime value but low engagement'). "
                    "Use the actual segment names and statistics from the analytics.\n\n"
                    "2. **Top Risk/Opportunity Drivers** — list the main predictors for each segment "
                    "from the ML model feature importance (e.g. contract type, usage frequency, "
                    "tenure, product adoption, payment method, etc.).\n\n"
                    "3. **Specific Campaign Ideas (4–6 ideas)** — give each campaign a concrete name "
                    "and mechanic directly tied to a finding from the data. Examples of mechanics:\n"
                    "   - Loyalty point programme for high-value retainable segments\n"
                    "   - Upgrade incentive for users on short-term or low-tier plans\n"
                    "   - Win-back flow for recently churned or disengaged users\n"
                    "   - Digital adoption bonus for users not using premium features\n"
                    "   - VIP early access for top-performing segments\n"
                    "   - Personalised discount for price-sensitive at-risk users\n"
                    "   Invent campaigns that directly address the actual drivers found in the data.\n\n"
                    "4. **Priority Channels** — recommend the best channel for each segment "
                    "(e.g. push notification, email, SMS, Instagram, LinkedIn, in-app message) "
                    "based on the segment profile.\n\n"
                    "5. **Measurable Targets** — specific, quantified goals "
                    "(e.g. 'Reduce Segment B churn from 42% to 28% in 90 days', "
                    "'Increase Segment A average spend by 15% in 60 days').\n\n"
                    + _analyst_section
                    + _scientist_section
                    + _chart_section
                    + f"## CDO Executive Summary:\n{banking_enhanced[:1500]}\n\n"
                    + "Write the brief in English. Be specific, data-driven, and creative. "
                    "Adapt all language to the industry and domain — do NOT assume banking. "
                    "This document is the primary input for the Digital Marketing creative team."
                )
                _ollama = ChatOllama(model="qwen3.5:cloud")
                _resp = _ollama.invoke([_HMsg(content=_handoff_prompt)])
                _raw_brief = _resp.content if hasattr(_resp, "content") else ""
                campaign_request = _strip_md_fences(_raw_brief)
                if campaign_request and len(campaign_request) > 200:
                    _push(task_id, "  Analyst → Marketing handoff brief generated ✓", "success")
                    derived_audience = "At-risk segments identified by analytics model (see brief)"
                    derived_goals    = "Drive retention and growth via data-driven targeted campaigns"
                else:
                    campaign_request = ""
        except Exception as _br_err:
            _push(task_id, f"  Campaign brief generation (non-critical): {_br_err}", "warn")

        # Fallback — if Gemini brief generation failed, build a structured fallback
        if not campaign_request or len(campaign_request) < 200:
            campaign_request = (
                f"Create a comprehensive digital marketing campaign for {_brand_name} "
                f"in the {req.industry} industry.\n\n"
                f"Campaign Type: {req.campaign_type} | Budget: {req.budget}\n\n"
                "Design specific, named campaign initiatives targeting the at-risk or high-opportunity "
                "segments identified in the analytics below. Include concrete campaign mechanics "
                "(loyalty programmes, upgrade incentives, win-back flows, digital adoption rewards). "
                "Each campaign should directly address a key driver found in the data.\n\n"
                f"Analytics Context:\n{banking_context[:3000]}"
            )
            _push(task_id, "  Using structured fallback campaign brief.", "warn")

        tasks[task_id]["campaign_brief"] = campaign_request

        # ── Phase 3a: Marketing LangGraph pre-planning (5 nodes, Gemini) ─────
        _push(task_id, "  Phase 3a — Marketing LangGraph: 5-node Gemini planning…", "info")
        from graphs.marketing_graph import run_marketing_analysis as _mkt_lg_run
        mkt_plan = _mkt_lg_run(
            task_description  = campaign_request,
            brand_name        = _brand_name,
            industry          = req.industry,
            target_audience   = derived_audience,
            campaign_goals    = derived_goals,
            budget            = req.budget,
            competitors       = req.competitors,
            campaign_type     = req.campaign_type,
            analytics_context = analytics_context,
        )
        tasks[task_id]["marketing_langgraph_plan"] = mkt_plan
        _push(task_id, "  Phase 3a complete — Marketing LangGraph planning done ✓", "success")

        # ── Phase 3b: Marketing CrewAI (4 agents) ────────────────────────────
        _push(task_id, "  Phase 3b — Marketing CrewAI: 4-agent sequential execution…", "info")
        from crew_marketing import run_marketing_crew
        marketing_result = run_marketing_crew(
            campaign_request  = campaign_request,
            brand_name        = _brand_name,
            industry          = req.industry,
            target_audience   = derived_audience,
            campaign_goals    = derived_goals,
            budget            = req.budget,
            competitors       = req.competitors,
            campaign_type     = req.campaign_type,
            analytics_context = analytics_context,
            langgraph_plan    = mkt_plan,
        )

        # Load marketing content sidecar
        marketing_content_items: list = []
        _content_file = BASE_DIR / "outputs" / "content" / "session_content.json"
        if _content_file.exists():
            try:
                all_entries = json.loads(_content_file.read_text(encoding="utf-8"))
                _seen_keys: set = set()
                for entry in all_entries:
                    key = f"{entry.get('type', '')}-{entry.get('title', '')}"
                    if key not in _seen_keys:
                        _seen_keys.add(key)
                        marketing_content_items.append(entry)
                _push(task_id, f"  Marketing content items loaded: {len(marketing_content_items)}", "info")
            except Exception as _e:
                _push(task_id, f"  Marketing content load failed: {_e}", "warn")

        # ── Fallback: generate content directly if crew produced nothing ───────
        if not marketing_content_items:
            _push(task_id, "  Content crew produced no items — running direct content fallback…", "warn")
            try:
                from tools.mkt_content import write_ad_copy, generate_social_posts, generate_content_report
                _fb_brand = req.brand_name or "Brand"
                _fb_campaign = req.campaign_type or "Retention"
                _fb_audience = req.target_audience or "General audience"

                _push(task_id, "  Fallback: generating ad copy…", "info")
                write_ad_copy(
                    brand=_fb_brand,
                    product=f"{_fb_brand} services",
                    audience=_fb_audience,
                    platform="Google",
                    tone="professional",
                )
                write_ad_copy(
                    brand=_fb_brand,
                    product=f"{_fb_brand} services",
                    audience=_fb_audience,
                    platform="Meta",
                    tone="inspirational",
                )

                _push(task_id, "  Fallback: generating social posts…", "info")
                generate_social_posts(
                    brand=_fb_brand,
                    campaign=f"{_fb_campaign} campaign for {_fb_audience}",
                    platforms="Instagram,Twitter,LinkedIn",
                    num_posts=3,
                )

                _push(task_id, "  Fallback: generating content report…", "info")
                generate_content_report(
                    summary=marketing_result[:1500] if marketing_result else f"{_fb_brand} {_fb_campaign} campaign content assets.",
                    brand_name=_fb_brand,
                )

                # Re-read the file
                if _content_file.exists():
                    all_entries = json.loads(_content_file.read_text(encoding="utf-8"))
                    _seen_keys2: set = set()
                    for entry in all_entries:
                        key = f"{entry.get('type', '')}-{entry.get('title', '')}"
                        if key not in _seen_keys2:
                            _seen_keys2.add(key)
                            marketing_content_items.append(entry)
                    _push(task_id, f"  Fallback content generated: {len(marketing_content_items)} items", "success")
            except Exception as _fb_err:
                _push(task_id, f"  Content fallback failed: {_fb_err}", "warn")

        # All downloads (banking + marketing reports + videos)
        all_reports = _list_files("outputs/reports", [".pdf", ".pptx", ".md"])
        all_videos  = [item for item in marketing_content_items if item.get("type") == "video"]

        tasks[task_id].update({
            "status":                  "completed",
            "completed_at":            datetime.now().isoformat(),
            "marketing_report":        marketing_result,
            "marketing_content_items": marketing_content_items,
            "all_reports":             all_reports,
            "all_videos":              all_videos,
        })

        _push(task_id, "━" * 62, "success")
        _push(task_id, "  All 3 phases complete! Banking + Marketing analysis done.", "success")
        _push(task_id, f"  Reports: {len(all_reports)}  Videos: {len(all_videos)}  Content items: {len(marketing_content_items)}", "success")
        _push(task_id, "━" * 62, "success")

    except Exception as exc:
        import traceback
        err_trace = traceback.format_exc()
        tasks[task_id].update({
            "status":       "error",
            "error":        str(exc),
            "completed_at": datetime.now().isoformat(),
        })
        _push(task_id, f"ERROR: {exc}", "error")
        _push(task_id, err_trace, "error")

    finally:
        sys.stdout = original_stdout
        done_type = "complete" if tasks[task_id]["status"] == "completed" else (
            "aborted" if tasks[task_id].get("hitl_aborted") else "error"
        )
        try:
            log_queues[task_id].put_nowait({
                "type": done_type, "task_id": task_id,
                "timestamp": datetime.now().isoformat()
            })
        except Exception:
            pass
        hitl_events.pop(task_id, None)


def _run_with_lock(task_id: str, req: CombinedRequest):
    with _analysis_lock:
        _run_analysis(task_id, req)


# ── Routes ────────────────────────────────────────────────────────────────────

@app.get("/", include_in_schema=False)
async def root():
    return FileResponse(str(BASE_DIR / "static" / "index.html"))


@app.post("/api/analyze")
async def start_analysis(req: CombinedRequest):
    """Start a new combined analysis run (background thread)."""
    if _analysis_lock.locked():
        raise HTTPException(status_code=409, detail="An analysis is already running.")

    task_id = str(uuid.uuid4())
    tasks[task_id] = {
        "status":                  "running",
        "started_at":              datetime.now().isoformat(),
        "completed_at":            None,
        "brand_name":              req.brand_name,
        "database_uri":            req.database_uri,
        "error":                   None,
        "logs":                    [],
        "langgraph_plan":            {},
        "marketing_langgraph_plan": {},
        "campaign_brief":           "",
        "hitl_aborted":            False,
        # Banking outputs
        "banking_report":          "",
        "banking_blog_sections":   [],
        "banking_charts":          [],
        # Marketing outputs
        "marketing_report":        "",
        "marketing_content_items": [],
        "all_reports":             [],
        "all_videos":              [],
        # Conversation
        "conversation":            [],
    }
    log_queues[task_id]  = queue.Queue()
    hitl_events[task_id] = threading.Event()

    thread = threading.Thread(target=_run_with_lock, args=(task_id, req), daemon=True)
    thread.start()

    return {"task_id": task_id, "status": "started"}


@app.get("/api/status/{task_id}")
async def get_status(task_id: str):
    if task_id not in tasks:
        raise HTTPException(status_code=404, detail="Task not found.")
    t = tasks[task_id]
    return {
        "task_id":      task_id,
        "status":       t["status"],
        "started_at":   t["started_at"],
        "completed_at": t.get("completed_at"),
        "error":        t.get("error"),
        "log_lines":    len(t.get("logs", [])),
    }


@app.get("/api/results/{task_id}")
async def get_results(task_id: str):
    if task_id not in tasks:
        raise HTTPException(status_code=404, detail="Task not found.")
    t = tasks[task_id]
    if t["status"] == "running":
        raise HTTPException(status_code=202, detail="Analysis still running.")
    if t["status"] == "error":
        raise HTTPException(status_code=500, detail=t.get("error", "Unknown error."))

    return {
        "task_id":                 task_id,
        "status":                  "completed",
        "banking_report":          t.get("banking_report", ""),
        "banking_blog_sections":   t.get("banking_blog_sections", []),
        "banking_charts":          t.get("banking_charts", []),
        "marketing_report":        t.get("marketing_report", ""),
        "campaign_brief":          t.get("campaign_brief", ""),
        "marketing_content_items": t.get("marketing_content_items", []),
        "all_reports":             t.get("all_reports", []),
        "all_videos":              t.get("all_videos", []),
        "langgraph_plan":          t.get("langgraph_plan", {}),
        "completed_at":            t.get("completed_at"),
    }


@app.get("/api/logs/{task_id}")
async def get_logs(task_id: str):
    if task_id not in tasks:
        raise HTTPException(status_code=404, detail="Task not found.")
    return {"task_id": task_id, "logs": tasks[task_id].get("logs", [])}


@app.get("/api/stream/{task_id}")
async def stream_events(task_id: str):
    if task_id not in tasks:
        raise HTTPException(status_code=404, detail="Task not found.")

    async def generator():
        q = log_queues[task_id]
        buffered = list(tasks[task_id].get("logs", []))
        replayed = len(buffered)
        for entry in buffered:
            yield f"data: {json.dumps(entry)}\n\n"

        drained = 0
        while drained < replayed:
            try:
                q.get_nowait()
                drained += 1
            except queue.Empty:
                break

        while True:
            try:
                msg = q.get_nowait()
                yield f"data: {json.dumps(msg)}\n\n"
                if msg.get("type") in ("complete", "error", "aborted"):
                    return
            except queue.Empty:
                if tasks[task_id]["status"] in ("completed", "error", "aborted"):
                    final = "complete" if tasks[task_id]["status"] == "completed" else "error"
                    yield f"data: {json.dumps({'type': final, 'task_id': task_id})}\n\n"
                    return
                await asyncio.sleep(0.35)

    return StreamingResponse(generator(), media_type="text/event-stream",
                              headers={"Cache-Control": "no-cache",
                                       "X-Accel-Buffering": "no",
                                       "Connection": "keep-alive"})


@app.get("/api/charts")
async def list_charts():
    return {"charts": _list_files("outputs/charts", [".png"])}


@app.get("/api/reports")
async def list_reports():
    return {"reports": _list_files("outputs/reports", [".pdf", ".pptx", ".md"])}


@app.get("/api/videos")
async def list_videos():
    return {"videos": _list_files("outputs/videos", [".mp4"])}


@app.get("/api/tasks")
async def list_tasks():
    return [
        {"task_id": tid, "status": t["status"],
         "started_at": t["started_at"], "completed_at": t.get("completed_at"),
         "brand_name": t.get("brand_name", "")}
        for tid, t in tasks.items()
    ]


@app.post("/api/approve/{task_id}")
async def approve_hitl(task_id: str, req: ApproveRequest = None):
    if req is None:
        req = ApproveRequest()
    if task_id not in tasks:
        raise HTTPException(status_code=404, detail="Task not found.")
    if tasks[task_id]["status"] != "awaiting_approval":
        raise HTTPException(status_code=400, detail="Task is not awaiting approval.")
    if task_id not in hitl_events:
        raise HTTPException(status_code=400, detail="No HITL event found for this task.")

    if req.abort:
        tasks[task_id]["hitl_aborted"] = True

    hitl_events[task_id].set()
    return {"task_id": task_id, "aborted": req.abort}


_SUPPORTED_EXTS = {
    ".csv", ".xlsx", ".xls",          # structured
    ".pdf", ".docx", ".doc",           # documents
    ".pptx", ".ppt",                   # presentations
}
_STRUCTURED_EXTS = {".csv", ".xlsx", ".xls"}


@app.post("/api/upload")
async def upload_data_files(files: List[UploadFile] = File(...)):
    """
    Upload one or more files for analysis.
    Supported: CSV, Excel (.csv/.xlsx/.xls) — structured data loaded into SQLite.
    Supported: PDF, Word, PowerPoint (.pdf/.docx/.pptx) — text extracted as document context.
    Multiple files may be uploaded at once.
    """
    if not files:
        raise HTTPException(status_code=400, detail="No files provided.")

    upload_dir = BASE_DIR / "outputs" / "uploads"
    upload_dir.mkdir(parents=True, exist_ok=True)

    results = []
    for file in files:
        if not file.filename:
            continue
        ext = Path(file.filename).suffix.lower()
        if ext not in _SUPPORTED_EXTS:
            raise HTTPException(
                status_code=400,
                detail=f"'{file.filename}': unsupported format. Allowed: CSV, Excel, PDF, Word (.docx), PowerPoint (.pptx)."
            )

        ts        = datetime.now().strftime("%Y%m%d_%H%M%S_%f")
        safe_name = "".join(c if c.isalnum() or c in "._-" else "_" for c in file.filename)
        raw_path  = upload_dir / f"{ts}_{safe_name}"
        content   = await file.read()
        raw_path.write_bytes(content)

        file_id = ts

        if ext in _STRUCTURED_EXTS:
            # ── Structured: load into SQLite ──────────────────────────────
            try:
                df = pd.read_csv(raw_path) if ext == ".csv" else pd.read_excel(raw_path)
            except Exception as e:
                raise HTTPException(status_code=422, detail=f"Could not read '{file.filename}': {e}")

            tbl = Path(file.filename).stem.lower()
            tbl = ("".join(c if c.isalnum() else "_" for c in tbl))[:50] or "data"

            sqlite_path = upload_dir / f"{ts}_data.db"
            try:
                from sqlalchemy import create_engine as _ce_up
                pd.DataFrame(df).to_sql(tbl, _ce_up(f"sqlite:///{sqlite_path}"),
                                        if_exists="replace", index=False)
            except Exception as e:
                raise HTTPException(status_code=500, detail=f"Could not store '{file.filename}': {e}")

            _uploaded_files[file_id] = {
                "original_name": file.filename,
                "file_type":     "structured",
                "format":        ext.lstrip("."),
                "table_name":    tbl,
                "sqlite_uri":    f"sqlite:///{sqlite_path}",
                "extracted_text": None,
                "rows":          len(df),
                "columns":       list(df.columns),
            }
            results.append({
                "file_id":    file_id,
                "filename":   file.filename,
                "type":       "structured",
                "format":     ext.lstrip("."),
                "table_name": tbl,
                "rows":       len(df),
                "columns":    list(df.columns),
            })

        else:
            # ── Document: extract text ────────────────────────────────────
            if ext == ".pdf":
                text = _extract_pdf_text(str(raw_path))
            elif ext in (".docx", ".doc"):
                text = _extract_word_text(str(raw_path))
            else:
                text = _extract_pptx_text(str(raw_path))

            word_count = len(text.split()) if text else 0
            _uploaded_files[file_id] = {
                "original_name":  file.filename,
                "file_type":      "document",
                "format":         ext.lstrip("."),
                "table_name":     None,
                "sqlite_uri":     None,
                "extracted_text": text[:60000],
                "rows":           None,
                "columns":        None,
            }
            results.append({
                "file_id":    file_id,
                "filename":   file.filename,
                "type":       "document",
                "format":     ext.lstrip("."),
                "word_count": word_count,
                "preview":    text[:200] if text else "",
            })

    return {"files": results, "count": len(results)}


# ── Agentic Chat Helpers ───────────────────────────────────────────────────

def _execute_viz_code(code: str, output_path: str) -> tuple:
    """
    Execute matplotlib/seaborn visualization code safely.
    Returns (success: bool, message: str).
    """
    BLOCKED = [
        "import os", "import sys", "import subprocess", "import socket",
        "import requests", "import pathlib", "__import__", "eval(", "exec(",
        "open(", "os.system", "os.popen", "shutil",
    ]
    for b in BLOCKED:
        if b in code:
            return False, f"Blocked pattern '{b}' not allowed in visualization code."

    setup = (
        "import matplotlib\nmatplotlib.use('Agg')\n"
        "import matplotlib.pyplot as plt\nimport pandas as pd\nimport numpy as np\n"
        "import json\nimport seaborn as sns\nfrom datetime import datetime\n"
        f"OUTPUT_PATH = {repr(output_path)}\n"
        "plt.style.use('dark_background')\n"
        "plt.rcParams.update({"
        "'figure.facecolor':'#0d1117','axes.facecolor':'#161b22',"
        "'text.color':'#e6edf3','axes.labelcolor':'#e6edf3',"
        "'xtick.color':'#8b949e','ytick.color':'#8b949e',"
        "'axes.edgecolor':'#30363d','grid.color':'#21262d'"
        "})\n"
    )
    full_code = setup + "\n" + code
    if "plt.show()" in full_code:
        full_code = full_code.replace("plt.show()", 'plt.savefig(OUTPUT_PATH, dpi=120, bbox_inches="tight"); plt.close()')
    if "plt.savefig" not in full_code:
        full_code += '\nplt.savefig(OUTPUT_PATH, dpi=120, bbox_inches="tight")\nplt.close()\n'
    try:
        exec(full_code, {})  # noqa: S102
        return True, "Chart generated successfully."
    except Exception as exc:
        return False, f"Code execution error: {exc}"


_VIZ_KEYWORDS = (
    "visualiz", "chart", "graph", "plot", "bar chart", "pie chart",
    "histogram", "scatter", "heatmap", "line chart", "show me a",
    "draw", "figure", "barchart", "linechart", "boxplot",
)


async def _agentic_chat(
    message: str,
    task_ctx: dict,
    history: list,
    task_id: str,
) -> tuple:
    """
    NVIDIA LLaMA agentic chat with intent-based visualization.

    Strategy:
      - Detect visualization intent directly from the user's message (keyword match).
      - If viz requested → two-pass: (1) force code generation, (2) get explanation.
      - Otherwise → single-pass conversational answer.

    Returns (answer: str, chart_urls: list[str], security_flags: list[str]).
    """
    from config import query_nvidia, nvidia_client
    if not nvidia_client:
        return "NVIDIA API key not configured — cannot process this request.", [], []

    banking_report   = (task_ctx.get("banking_report") or "")[:4000]
    marketing_report = (task_ctx.get("marketing_report") or "")[:1500]
    chart_analyses   = "\n\n".join(
        f"**{s['title']}**:\n{s.get('gemini_analysis', '')}"
        for s in (task_ctx.get("banking_blog_sections") or [])
        if s.get("gemini_analysis")
    )[:2000]

    context_block = (
        "## Analytics Report:\n" + banking_report + "\n\n"
        + (f"## Chart Analyses:\n{chart_analyses}\n\n" if chart_analyses else "")
        + "## Marketing Campaign Brief:\n" + marketing_report
    )

    # ── Detect visualization intent ───────────────────────────────────────────
    msg_lower = message.lower()
    wants_viz = any(kw in msg_lower for kw in _VIZ_KEYWORDS)

    chart_urls: list = []

    if wants_viz:
        # ── Pass 1: Force Python code generation ──────────────────────────────
        code_system = (
            "You are a Python data visualization expert.\n"
            "Your ONLY job is to write executable matplotlib/seaborn code.\n"
            "Output a single ```python ... ``` code block and nothing else.\n\n"
            "STRICT RULES:\n"
            "1. Use ONLY: matplotlib, seaborn, pandas, numpy, json, datetime\n"
            "2. Do NOT import os, sys, subprocess, requests, pathlib or any other module\n"
            "3. Save with exactly: plt.savefig(OUTPUT_PATH, dpi=150, bbox_inches='tight')\n"
            "4. Do NOT call plt.show()\n"
            "5. Use hardcoded/synthetic data that matches the analytics report context\n"
            "6. Chart must be dark-themed (plt.style.use('dark_background') is already applied)\n"
            "7. Add a clear title and axis labels"
        )
        code_user = (
            f"{context_block}\n\n"
            f"User request: {message}\n\n"
            "Write the Python visualization code now. Output ONLY the ```python``` block."
        )
        code_messages = [
            {"role": "system", "content": code_system},
            {"role": "user",   "content": code_user},
        ]
        try:
            code_response = await asyncio.to_thread(query_nvidia, code_messages, 0.1, 1800)
        except Exception as exc:
            return f"NVIDIA LLaMA error generating visualization: {exc}", [], []

        code_match = re.search(r"```python\s*([\s\S]*?)```", code_response, re.IGNORECASE)
        if not code_match:
            # Model didn't wrap in fences — try treating the whole response as code
            code_match = re.search(r"(import matplotlib[\s\S]*)", code_response)

        if code_match:
            code = code_match.group(1).strip()
            ts_viz         = datetime.now().strftime("%Y%m%d_%H%M%S_%f")
            chart_filename = f"chat_viz_{task_id[:8]}_{ts_viz}.png"
            chart_path     = str(BASE_DIR / "outputs" / "charts" / chart_filename)
            chart_url      = f"/outputs/charts/{chart_filename}"

            success, exec_msg = _execute_viz_code(code, chart_path)
            if success and Path(chart_path).exists():
                chart_urls.append(chart_url)
                viz_status = f"Chart successfully generated and saved at {chart_url}."
            else:
                viz_status = f"Visualization failed: {exec_msg}"
                print(f"[Chat] Viz exec failed: {exec_msg}\nCode:\n{code[:400]}", flush=True)
        else:
            viz_status = "Could not extract Python code from the model response."
            print(f"[Chat] No code block found in: {code_response[:300]}", flush=True)

        # ── Pass 2: Get natural-language explanation ───────────────────────────
        explain_system = (
            "You are an expert data analyst assistant.\n"
            "Answer based on the analytics context provided.\n"
            "Be concise and insightful. Do NOT generate code."
        )
        explain_user = (
            f"{context_block}\n\n"
            f"User asked: {message}\n\n"
            f"Visualization status: {viz_status}\n\n"
            "Provide a concise explanation of what the chart shows and key insights from it. "
            "If the chart failed, explain what it would have shown and the key data points."
        )
        explain_messages = [{"role": "system", "content": explain_system}]
        for turn in history[-4:]:
            explain_messages.append({"role": "user",      "content": turn["user"]})
            explain_messages.append({"role": "assistant", "content": turn["assistant"]})
        explain_messages.append({"role": "user", "content": explain_user})

        try:
            answer = await asyncio.to_thread(query_nvidia, explain_messages, 0.3, 1024)
        except Exception as exc:
            answer = f"Chart generated. {viz_status}"

        return answer, chart_urls, []

    # ── Conversational answer (no visualization needed) ───────────────────────
    conv_system = (
        "You are an expert data analyst AI assistant.\n"
        "Answer questions based strictly on the analytics results provided.\n"
        "Be concise, specific, and reference actual data from the report.\n"
        "Do NOT generate code or charts — just answer the question directly.\n\n"
        + context_block
    )
    conv_messages = [{"role": "system", "content": conv_system}]
    for turn in history[-6:]:
        conv_messages.append({"role": "user",      "content": turn["user"]})
        conv_messages.append({"role": "assistant", "content": turn["assistant"]})
    conv_messages.append({"role": "user", "content": message})

    try:
        answer = await asyncio.to_thread(query_nvidia, conv_messages, 0.3, 1500)
    except Exception as exc:
        return f"NVIDIA LLaMA error: {exc}", [], []

    return answer, [], []


@app.post("/api/chat/{task_id}")
async def chat(task_id: str, req: ChatRequest):
    """Agentic chat powered by NVIDIA LLaMA with on-demand visualization tool."""
    if task_id not in tasks:
        raise HTTPException(status_code=404, detail="Task not found.")
    t = tasks[task_id]
    if t["status"] != "completed":
        raise HTTPException(status_code=400, detail="Analysis not yet completed.")

    message = req.message.strip()
    if len(message) > 2000:
        raise HTTPException(status_code=400, detail="Message too long (max 2000 characters).")
    if _is_sql_injection(message):
        _audit_log({"timestamp": datetime.now().isoformat(), "task_id": task_id,
                    "flags": ["sql_injection"], "blocked": True})
        raise HTTPException(status_code=400, detail="Blocked: SQL injection pattern detected.")
    if _is_guardrail_violation(message):
        _audit_log({"timestamp": datetime.now().isoformat(), "task_id": task_id,
                    "flags": ["guardrail_violation"], "blocked": True})
        raise HTTPException(status_code=400, detail="Blocked: out of scope or policy violation.")

    clean_message, pii_in_input = _redact_pii(message)

    answer, chart_urls, _ = await _agentic_chat(
        message   = clean_message,
        task_ctx  = t,
        history   = t.get("conversation", [])[-10:],
        task_id   = task_id,
    )

    reply, pii_in_output = _redact_pii(answer)
    security_flags = []
    if pii_in_input:  security_flags.append("pii_redacted_input")
    if pii_in_output: security_flags.append("pii_redacted_output")

    _audit_log({"timestamp": datetime.now().isoformat(), "task_id": task_id,
                "input_len": len(message), "flags": security_flags, "blocked": False,
                "output_len": len(reply), "chart_urls": chart_urls})

    t["conversation"].append({"user": clean_message, "assistant": reply,
                               "chart_urls": chart_urls})
    return {"reply": reply, "chart_urls": chart_urls, "security_flags": security_flags}


@app.get("/api/chat/history/{task_id}")
async def chat_history(task_id: str):
    if task_id not in tasks:
        raise HTTPException(status_code=404, detail="Task not found.")
    return {"history": tasks[task_id].get("conversation", [])}


if __name__ == "__main__":
    import uvicorn
    uvicorn.run("app:app", host="127.0.0.1", port=8002, reload=False)
