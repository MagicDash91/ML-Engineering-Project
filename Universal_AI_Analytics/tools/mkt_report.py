"""
tools/mkt_report.py – Marketing report generation tools (PDF, PPTX, Markdown).
Mirrors Digital_Marketing_Agent/tools/report_tools.py with paths adjusted.
"""
from __future__ import annotations

import os
import sys
import re
import ast
import json
from datetime import datetime
from pathlib import Path

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from crewai.tools import tool

_BASE       = Path(__file__).parent.parent / "outputs"
_REPORTS    = _BASE / "reports"
_REPORTS.mkdir(parents=True, exist_ok=True)


def _gemini(prompt: str, system: str = "", _timeout: int = 130) -> str:
    """Call Ollama qwen3.5:cloud with hard wall-clock timeout."""
    import concurrent.futures

    def _call() -> str:
        import time
        from langchain_community.chat_models import ChatOllama
        from langchain_core.messages import HumanMessage, SystemMessage
        msgs = []
        if system:
            msgs.append(SystemMessage(content=system))
        msgs.append(HumanMessage(content=prompt))
        last_exc = None
        for attempt in range(3):
            try:
                llm = ChatOllama(model="qwen3.5:cloud")
                resp = llm.invoke(msgs)
                return resp.content if hasattr(resp, "content") else str(resp)
            except Exception as exc:
                last_exc = exc
                if attempt < 2:
                    print(f"[MktReport] Ollama error (attempt {attempt+1}/3): {exc} — retrying in 5s", flush=True)
                    time.sleep(5)
        return f"Unable to generate response: {last_exc}"

    ex = concurrent.futures.ThreadPoolExecutor(max_workers=1)
    fut = ex.submit(_call)
    try:
        result = fut.result(timeout=_timeout)
        ex.shutdown(wait=False)
        return result
    except concurrent.futures.TimeoutError:
        ex.shutdown(wait=False)
        return "Response timed out."


@tool("generate_text_report")
def generate_text_report(content: str, brand_name: str) -> str:
    """
    Save a markdown campaign report.

    Args:
        content:    Report content (markdown).
        brand_name: Brand name for filename.

    Returns:
        Path to saved markdown file.
    """
    ts   = datetime.now().strftime("%Y%m%d_%H%M%S")
    safe = "".join(c if c.isalnum() else "_" for c in brand_name)
    path = _REPORTS / f"campaign_report_{safe}_{ts}.md"
    path.write_text(content, encoding="utf-8")
    print(f"[MktReport] Markdown report saved → {path}", flush=True)
    return f"Report saved to {path}"


@tool("generate_pdf_report")
def generate_pdf_report(content: str, brand_name: str) -> str:
    """
    Generate a professional PDF marketing campaign report.

    Args:
        content:    Report content (markdown).
        brand_name: Brand name for the report.

    Returns:
        Path to saved PDF file.
    """
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import inch, cm
        from reportlab.lib import colors
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, HRFlowable
        from reportlab.lib.enums import TA_CENTER, TA_JUSTIFY

        ts       = datetime.now().strftime("%Y%m%d_%H%M%S")
        safe     = "".join(c if c.isalnum() else "_" for c in brand_name)
        pdf_path = str(_REPORTS / f"campaign_report_{safe}_{ts}.pdf")

        doc  = SimpleDocTemplate(pdf_path, pagesize=A4,
                                  rightMargin=2*cm, leftMargin=2*cm,
                                  topMargin=2.5*cm, bottomMargin=2*cm)
        base = getSampleStyleSheet()
        GREEN      = colors.HexColor("#2E7D32")
        LIGHT_GREEN= colors.HexColor("#388E3C")

        s_title = ParagraphStyle("s_title", parent=base["Title"], fontSize=22, textColor=GREEN,
                                  spaceAfter=18, alignment=TA_CENTER, fontName="Helvetica-Bold")
        s_h1    = ParagraphStyle("s_h1", parent=base["Heading1"], fontSize=14, textColor=GREEN,
                                  spaceBefore=16, spaceAfter=8, fontName="Helvetica-Bold")
        s_h2    = ParagraphStyle("s_h2", parent=base["Heading2"], fontSize=12, textColor=LIGHT_GREEN,
                                  spaceBefore=10, spaceAfter=5, fontName="Helvetica-Bold")
        s_body  = ParagraphStyle("s_body", parent=base["Normal"], fontSize=10, leading=14,
                                  spaceAfter=7, alignment=TA_JUSTIFY, fontName="Helvetica")
        s_sub   = ParagraphStyle("s_sub", parent=base["Normal"], fontSize=10, alignment=TA_CENTER,
                                  textColor=colors.grey)

        report_title = f"{brand_name} — Digital Marketing Campaign Report"
        story = []
        story.append(Spacer(1, 0.6*inch))
        story.append(Paragraph(report_title, s_title))
        story.append(Paragraph(f"Generated: {datetime.now().strftime('%B %d, %Y  %H:%M')}", s_sub))
        story.append(HRFlowable(width="100%", thickness=2, color=GREEN))
        story.append(Spacer(1, 0.4*inch))

        _bold_re = re.compile(r"\*\*(.+?)\*\*")
        for raw_line in content.split("\n"):
            line = raw_line.strip()
            if not line:
                story.append(Spacer(1, 0.08*inch))
            elif line.startswith("# "):
                story.append(Paragraph(line[2:], s_title))
            elif line.startswith("## "):
                story.append(Paragraph(line[3:], s_h1))
            elif line.startswith("### "):
                story.append(Paragraph(line[4:], s_h2))
            elif line.startswith(("- ", "* ", "• ")):
                body = _bold_re.sub(r"<b>\1</b>", line[2:])
                story.append(Paragraph(f"&bull;&nbsp;{body}", s_body))
            else:
                story.append(Paragraph(_bold_re.sub(r"<b>\1</b>", line), s_body))

        doc.build(story)
        print(f"[MktReport] PDF saved → {pdf_path}", flush=True)
        return f"PDF report saved to {pdf_path}"
    except Exception as e:
        return f"[PDF generation failed: {e}]"


@tool("generate_ppt_report")
def generate_ppt_report(content: str, brand_name: str) -> str:
    """
    Generate a widescreen (16:9) PowerPoint marketing campaign presentation.

    Args:
        content:    Report content (markdown).
        brand_name: Brand name for the report.

    Returns:
        Path to saved PPTX file.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        prs = Presentation()
        prs.slide_width  = Inches(13.33)
        prs.slide_height = Inches(7.5)

        C_PRIMARY = RGBColor(0x2E, 0x7D, 0x32)   # green
        C_ACCENT  = RGBColor(0x38, 0x8E, 0x3C)
        C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
        C_DARK    = RGBColor(0x37, 0x47, 0x4F)
        C_SUBTEXT = RGBColor(0xC8, 0xE6, 0xC9)
        BLANK     = prs.slide_layouts[6]

        def _rect(slide, l, t, w, h, color):
            shp = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
            shp.fill.solid()
            shp.fill.fore_color.rgb = color
            shp.line.fill.background()

        def _textbox(slide, l, t, w, h, text, size, bold=False, color=C_WHITE, align=PP_ALIGN.LEFT):
            tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = text
            p.alignment = align
            run = p.runs[0] if p.runs else p.add_run()
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color

        def add_cover(title, subtitle):
            slide = prs.slides.add_slide(BLANK)
            bg = slide.background.fill
            bg.solid()
            bg.fore_color.rgb = C_PRIMARY
            _textbox(slide, 1.0, 2.3, 11.33, 1.6, title, 36, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
            _textbox(slide, 1.0, 4.1, 11.33, 0.9, subtitle, 18, color=C_SUBTEXT, align=PP_ALIGN.CENTER)

        def add_content(section_title, lines):
            slide = prs.slides.add_slide(BLANK)
            _rect(slide, 0, 0, 13.33, 1.15, C_PRIMARY)
            _textbox(slide, 0.3, 0.15, 12.73, 0.85, section_title, 22, bold=True)
            tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.35), Inches(12.33), Inches(5.7))
            tf = tb.text_frame
            tf.word_wrap = True
            for i, ln in enumerate(lines[:16]):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.space_before = Pt(5)
                cleaned = ln.replace("**", "").replace("##", "").replace("#", "").strip()
                p.text = cleaned
                run = p.runs[0] if p.runs else p.add_run()
                run.font.size = Pt(12)
                run.font.bold = ln.startswith("**")
                run.font.color.rgb = C_DARK

        add_cover(f"{brand_name} — Digital Marketing Campaign",
                  f"Campaign Report  •  {datetime.now().strftime('%B %Y')}")

        import re as _re
        _normalised = _re.sub(r'\n#{2,3}\s+', '\n## ', content)
        sections = _normalised.split("\n## ")
        for sec in sections[:8]:
            raw_lines = sec.strip().split("\n")
            sec_title = raw_lines[0].replace("#", "").strip()
            body      = [l.strip() for l in raw_lines[1:] if l.strip()]
            if sec_title and body:
                add_content(sec_title, body)

        add_cover("Thank You", "Digital Marketing Campaign — Summary Complete")

        ts       = datetime.now().strftime("%Y%m%d_%H%M%S")
        safe     = "".join(c if c.isalnum() else "_" for c in brand_name)
        ppt_path = str(_REPORTS / f"marketing_presentation_{safe}_{ts}.pptx")
        prs.save(ppt_path)
        print(f"[MktReport] PPTX saved → {ppt_path}", flush=True)
        return f"PowerPoint saved to {ppt_path}"
    except Exception as e:
        return f"[PPTX generation failed: {e}]"
