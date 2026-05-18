import os
import sys
import uuid
import json
import asyncio
from datetime import datetime
from pathlib import Path

os.environ["OTEL_SDK_DISABLED"]         = "true"
os.environ["CREWAI_TELEMETRY_OPT_OUT"] = "true"

_BASE = Path(__file__).parent
sys.path.insert(0, str(_BASE))

from fastapi import FastAPI, HTTPException, File, UploadFile
from fastapi.responses import HTMLResponse, StreamingResponse
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel
from typing import List

import config  # noqa — side effects: loads .env, sets env vars

from meta_planner  import generate_plan, format_plan_for_display
from dynamic_graph import build_and_run

app = FastAPI(title="AGI System", version="1.0")

outputs_dir = _BASE / "outputs"
outputs_dir.mkdir(exist_ok=True)
(_BASE / "outputs" / "uploads").mkdir(parents=True, exist_ok=True)
app.mount("/outputs", StaticFiles(directory=str(outputs_dir)), name="outputs")

# Also serve Universal_AI_Analytics outputs so chart images render in the browser
_ua_outputs = _BASE.parent / "Universal_AI_Analytics" / "outputs"
if _ua_outputs.exists():
    app.mount("/ua-outputs", StaticFiles(directory=str(_ua_outputs)), name="ua-outputs")

# In-memory stores
tasks: dict          = {}
_uploaded_files: dict = {}

_STRUCTURED_EXTS = {".csv", ".xlsx", ".xls"}
_DOCUMENT_EXTS   = {".pdf", ".docx", ".doc", ".pptx", ".ppt"}
_SUPPORTED_EXTS  = _STRUCTURED_EXTS | _DOCUMENT_EXTS


# ── Text extraction helpers ───────────────────────────────────────────────────

def _extract_pdf_text(path: str) -> str:
    try:
        import pdfplumber
        with pdfplumber.open(path) as pdf:
            return "\n\n".join(p.extract_text() or "" for p in pdf.pages)
    except Exception:
        pass
    try:
        from pypdf import PdfReader
        r = PdfReader(path)
        return "\n\n".join(p.extract_text() or "" for p in r.pages)
    except Exception as e:
        return f"[PDF extraction failed: {e}]"


def _extract_word_text(path: str) -> str:
    try:
        from docx import Document
        doc = Document(path)
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except Exception as e:
        return f"[Word extraction failed: {e}]"


def _extract_pptx_text(path: str) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(path)
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text)
        return "\n\n".join(parts)
    except Exception as e:
        return f"[PPTX extraction failed: {e}]"


# ── Request models ────────────────────────────────────────────────────────────

class SolveRequest(BaseModel):
    question:              str
    db_uri:                str        = ""
    uploaded_file_ids:     List[str]  = []
    conversation_history:  List[dict] = []  # [{role, content}, ...]


def _classify_question(question: str) -> str:
    """Returns 'simple' or 'complex'. Simple = conversational/no tools. Complex = data/ML/research."""
    from config import query_ollama
    try:
        raw = query_ollama(
            [{
                "role": "user",
                "content": (
                    "Classify this question as SIMPLE or COMPLEX.\n\n"
                    "SIMPLE: conversational, follow-up, opinion, general knowledge, marketing ideas, "
                    "explanations — can be answered without databases, ML models, or web research.\n"
                    "COMPLEX: needs data analysis, SQL queries, machine learning, web research, "
                    "report generation, or multi-step investigation.\n\n"
                    f"Question: {question}\n\n"
                    "Reply with ONLY one word: simple or complex"
                ),
            }],
            temperature=0,
            max_tokens=5,
        )
        return "simple" if "simple" in raw.lower() else "complex"
    except Exception:
        return "complex"


# ── Background pipeline ───────────────────────────────────────────────────────

async def _run_pipeline(task_id: str, question: str, db_uri: str, file_ids: List[str],
                        conversation_history: List[dict]):
    t = tasks[task_id]
    q = t["queue"]

    async def emit(event_type, message, payload=None):
        line = json.dumps({"type": event_type, "message": message, "payload": payload or {}})
        await q.put(line)
        t["logs"].append({"type": event_type, "message": message, "ts": datetime.now().isoformat()})

    try:
        # ── Resolve uploaded files ─────────────────────────────────────────────
        document_context = ""
        active_db_uri    = db_uri or None

        # Clear any stale DB URI from a previous pipeline run — must not leak across requests
        if not active_db_uri:
            os.environ.pop("ACTIVE_DB_URI", None)

        if file_ids:
            doc_parts   = []
            sqlite_uris = []

            for fid in file_ids:
                info = _uploaded_files.get(fid)
                if not info:
                    continue
                if info["file_type"] == "structured" and info.get("sqlite_uri"):
                    sqlite_uris.append(info["sqlite_uri"])
                    await emit("info", f"📊 Loaded structured file: {info['original_name']} ({info.get('rows',0)} rows)")
                elif info["file_type"] == "document" and info.get("extracted_text"):
                    doc_parts.append(f"=== {info['original_name']} ===\n{info['extracted_text'][:8000]}")
                    await emit("info", f"📄 Loaded document: {info['original_name']}")

            # Merge multiple SQLite DBs or pick first
            if sqlite_uris and not active_db_uri:
                if len(sqlite_uris) == 1:
                    active_db_uri = sqlite_uris[0]
                else:
                    # Merge into one SQLite
                    import pandas as pd
                    from sqlalchemy import create_engine
                    ts_mrg   = datetime.now().strftime("%Y%m%d_%H%M%S")
                    mrg_path = _BASE / "outputs" / "uploads" / f"merged_{task_id[:8]}_{ts_mrg}.db"
                    from sqlalchemy import inspect as _sa_inspect
                    dest_eng = create_engine(f"sqlite:///{mrg_path}")
                    for uri in sqlite_uris:
                        src_eng = create_engine(uri)
                        for tbl in _sa_inspect(src_eng).get_table_names():
                            pd.read_sql_table(tbl, src_eng).to_sql(tbl, dest_eng, if_exists="replace", index=False)
                    active_db_uri = f"sqlite:///{mrg_path}"
                    await emit("info", "🔗 Merged structured files into one database")

            if doc_parts:
                document_context = "\n\n".join(doc_parts)

        # ── Classify question complexity ───────────────────────────────────────
        complexity = await asyncio.to_thread(_classify_question, question)

        if complexity == "simple":
            t["status"] = "running"
            await emit("info", "💬 Simple question detected — answering directly...")
            from config import query_ollama
            history_ctx = (conversation_history or [])[-6:]  # last 3 exchanges
            messages = [
                {
                    "role": "system",
                    "content": (
                        "You are a helpful AI analyst and strategist. "
                        "Answer clearly and concisely using markdown formatting. "
                        "Use headers, bullet points, and bold text where appropriate."
                    ),
                },
                *history_ctx,
                {"role": "user", "content": question},
            ]
            answer = await asyncio.to_thread(query_ollama, messages, 0.3, 2048)
            t["final_answer"] = answer
            t["status"]       = "completed"
            await emit("done", "✅ Done", {"final_answer": answer})
            return

        t["status"] = "planning"
        await emit("info", "🧠 Meta-planner analyzing your question...")

        ctx = {"db_uri": active_db_uri, "has_files": bool(file_ids)}
        if document_context:
            ctx["document_context"] = document_context[:3000]

        plan = await asyncio.to_thread(generate_plan, question, ctx, document_context)
        t["plan"]   = plan
        t["status"] = "running"

        plan_display = format_plan_for_display(plan)
        await emit("plan", f"📋 Plan ready — {len(plan['steps'])} steps",
                   {"plan": plan, "display": plan_display})

        inner_q = asyncio.Queue()

        async def forward_inner():
            while True:
                item = await inner_q.get()
                if item is None:
                    break
                line = json.dumps(item)
                await q.put(line)
                t["logs"].append({"type": item["type"], "message": item["message"], "ts": datetime.now().isoformat()})

        forwarder = asyncio.create_task(forward_inner())

        result = await build_and_run(
            user_question    = question,
            plan             = plan,
            db_uri           = active_db_uri,
            document_context = document_context,
            log_queue        = inner_q,
        )

        await inner_q.put(None)
        await forwarder

        t["step_results"]  = result["step_results"]
        t["final_answer"]  = result["final_answer"]
        t["status"]        = "completed"

        await emit("done", "✅ Pipeline complete", {"final_answer": result["final_answer"]})

    except Exception as e:
        t["status"] = "error"
        t["error"]  = str(e)
        await emit("error", f"❌ Pipeline error: {e}", {})
        print(f"[AGI] Pipeline error for {task_id}: {e}", flush=True)

    finally:
        await q.put(None)


# ── Endpoints ─────────────────────────────────────────────────────────────────

@app.get("/", response_class=HTMLResponse)
async def index():
    return HTMLResponse((_BASE / "static" / "index.html").read_text(encoding="utf-8"))


@app.post("/api/upload")
async def upload_files(files: List[UploadFile] = File(...)):
    if not files:
        raise HTTPException(status_code=400, detail="No files provided.")

    upload_dir = _BASE / "outputs" / "uploads"
    upload_dir.mkdir(parents=True, exist_ok=True)
    results = []

    for file in files:
        if not file.filename:
            continue
        ext = Path(file.filename).suffix.lower()
        if ext not in _SUPPORTED_EXTS:
            raise HTTPException(status_code=400,
                detail=f"'{file.filename}': unsupported. Allowed: CSV, Excel, PDF, Word, PowerPoint.")

        ts        = datetime.now().strftime("%Y%m%d_%H%M%S_%f")
        safe_name = "".join(c if c.isalnum() or c in "._-" else "_" for c in file.filename)
        raw_path  = upload_dir / f"{ts}_{safe_name}"
        content   = await file.read()
        raw_path.write_bytes(content)
        file_id   = ts

        if ext in _STRUCTURED_EXTS:
            import pandas as pd
            from sqlalchemy import create_engine
            try:
                df = pd.read_csv(raw_path) if ext == ".csv" else pd.read_excel(raw_path)
            except Exception as e:
                raise HTTPException(status_code=422, detail=f"Could not read '{file.filename}': {e}")

            tbl = ("".join(c if c.isalnum() else "_" for c in Path(file.filename).stem.lower()))[:50] or "data"
            sqlite_path = upload_dir / f"{ts}_data.db"
            df.to_sql(tbl, create_engine(f"sqlite:///{sqlite_path}"), if_exists="replace", index=False)

            _uploaded_files[file_id] = {
                "original_name": file.filename, "file_type": "structured",
                "format": ext.lstrip("."), "table_name": tbl,
                "sqlite_uri": f"sqlite:///{sqlite_path}",
                "extracted_text": None, "rows": len(df), "columns": list(df.columns),
            }
            results.append({"file_id": file_id, "filename": file.filename,
                            "type": "structured", "table_name": tbl,
                            "rows": len(df), "columns": list(df.columns)})
        else:
            if ext == ".pdf":
                text = _extract_pdf_text(str(raw_path))
            elif ext in (".docx", ".doc"):
                text = _extract_word_text(str(raw_path))
            else:
                text = _extract_pptx_text(str(raw_path))

            _uploaded_files[file_id] = {
                "original_name": file.filename, "file_type": "document",
                "format": ext.lstrip("."), "table_name": None,
                "sqlite_uri": None, "extracted_text": text[:60000],
                "rows": None, "columns": None,
            }
            results.append({"file_id": file_id, "filename": file.filename,
                            "type": "document", "word_count": len(text.split()),
                            "preview": text[:200]})

    return {"files": results, "count": len(results)}


@app.post("/api/solve")
async def solve(req: SolveRequest):
    if not req.question.strip():
        raise HTTPException(status_code=400, detail="Question cannot be empty.")

    task_id = str(uuid.uuid4())
    tasks[task_id] = {
        "task_id":      task_id,
        "question":     req.question,
        "db_uri":       req.db_uri,
        "status":       "queued",
        "plan":         None,
        "step_results": [],
        "final_answer": "",
        "error":        "",
        "logs":         [],
        "queue":        asyncio.Queue(),
        "created_at":   datetime.now().isoformat(),
    }

    asyncio.create_task(_run_pipeline(task_id, req.question, req.db_uri, req.uploaded_file_ids, req.conversation_history))
    return {"task_id": task_id}


@app.get("/api/stream/{task_id}")
async def stream(task_id: str):
    if task_id not in tasks:
        raise HTTPException(status_code=404, detail="Task not found.")
    q = tasks[task_id]["queue"]

    async def event_generator():
        last_ka = asyncio.get_event_loop().time()
        while True:
            try:
                line = await asyncio.wait_for(q.get(), timeout=1.0)
            except asyncio.TimeoutError:
                if asyncio.get_event_loop().time() - last_ka > 15:
                    yield ": keepalive\n\n"
                    last_ka = asyncio.get_event_loop().time()
                continue
            if line is None:
                yield "data: {\"type\":\"stream_end\"}\n\n"
                break
            yield f"data: {line}\n\n"
            last_ka = asyncio.get_event_loop().time()

    return StreamingResponse(event_generator(),
        media_type="text/event-stream",
        headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"})


@app.get("/api/status/{task_id}")
async def status(task_id: str):
    if task_id not in tasks:
        raise HTTPException(status_code=404, detail="Task not found.")
    t = tasks[task_id]
    return {"task_id": task_id, "status": t["status"], "question": t["question"],
            "plan": t["plan"], "step_results": t["step_results"],
            "final_answer": t["final_answer"], "error": t["error"]}


@app.get("/api/tasks")
async def list_tasks():
    return [{"task_id": tid, "question": t["question"][:80], "status": t["status"],
             "created_at": t["created_at"], "steps": len(t["step_results"])}
            for tid, t in tasks.items()]


@app.get("/api/logs/{task_id}")
async def get_logs(task_id: str):
    if task_id not in tasks:
        raise HTTPException(status_code=404, detail="Task not found.")
    return tasks[task_id]["logs"]
