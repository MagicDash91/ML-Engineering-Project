"""
app.py – AI Financial Risk & Compliance Monitor FastAPI Backend
================================================================

Two-phase pipeline:
  Phase 1 → LangGraph risk planning (5 nodes, Ollama)
  HITL    → Human approval checkpoint
  Phase 2 → Risk Analytics CrewAI (5 agents: Engineer, Scientist,
             LabelEncoder, Risk Analyst, CRO)
             + Ollama enrichment of the risk report

Inputs (at least one data source required alongside the question):
  - Database URI  → any SQLAlchemy-compatible URI
  - File upload   → CSV/Excel (→ SQLite), PDF/Word/PPTX (→ text context)

5-tab Bootstrap 5 dark UI:
  Console | Risk Report | Charts | Compliance | Downloads | Conversation

Run:
    python main.py   →  http://localhost:8003
"""

import os
import sys

os.environ["OTEL_SDK_DISABLED"]         = "true"
os.environ["CREWAI_TELEMETRY_OPT_OUT"] = "true"

import re
import json
import uuid
import queue
import asyncio
import threading
from datetime import datetime
from pathlib import Path
from typing import List, Optional
import pandas as pd

BASE_DIR = Path(__file__).parent
sys.path.insert(0, str(BASE_DIR))

for _d in ("outputs/charts", "outputs/reports", "outputs/models",
           "outputs/audit", "outputs/uploads", "static"):
    (BASE_DIR / _d).mkdir(parents=True, exist_ok=True)

AUDIT_DIR = BASE_DIR / "outputs" / "audit"

from fastapi import FastAPI, HTTPException, File, UploadFile
from fastapi.responses import StreamingResponse, FileResponse
from fastapi.staticfiles import StaticFiles
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel

app = FastAPI(
    title="AI Financial Risk & Compliance Monitor",
    description="Connect any database or upload files — AI agents assess, model, and report financial risk",
    version="1.0.0",
)

app.add_middleware(CORSMiddleware, allow_origins=["*"], allow_credentials=True,
                   allow_methods=["*"], allow_headers=["*"])

app.mount("/static",  StaticFiles(directory=str(BASE_DIR / "static")),  name="static")
app.mount("/outputs", StaticFiles(directory=str(BASE_DIR / "outputs")), name="outputs")

# ── In-memory store ───────────────────────────────────────────────────────────
tasks: dict        = {}
log_queues: dict   = {}
hitl_events: dict  = {}
_analysis_lock     = threading.Lock()
_uploaded_files: dict = {}

# ── Security patterns ─────────────────────────────────────────────────────────
_PII_PATTERNS = [
    (re.compile(r'\b[A-Za-z0-9._%+\-]+@[A-Za-z0-9.\-]+\.[A-Za-z]{2,}\b'), '[EMAIL]'),
    (re.compile(r'(\+62|62|0)[\s.\-]?8[1-9][0-9]{6,10}\b'), '[PHONE]'),
    (re.compile(r'\b\d{16}\b'), '[NIK/ACCT]'),
    (re.compile(r'\b(?:\d{4}[\s\-]?){3}\d{4}\b'), '[CARD]'),
    (re.compile(r'\b\d{10,15}\b'), '[ACCT]'),
]

_SQL_RE = re.compile(
    r'(union\s+(all\s+)?select)|(drop\s+(table|database))'
    r'|(insert\s+into\s+\w)|(delete\s+from\s+\w)|(update\s+\w+\s+set\s)'
    r'|(exec\s*\(|xp_cmdshell)|(\bor\b\s+[\'"]?\d[\'"]?\s*=\s*[\'"]?\d)'
    r'|(--\s*[\r\n]|;\s*--)|(\/\*[\s\S]*?\*\/)',
    re.IGNORECASE,
)

_GUARDRAIL_RE = re.compile(
    r'\b(ignore|forget|disregard)\s+(previous|all|above|your)\s+(instruction|rule|prompt)'
    r'|\b(jailbreak|prompt.injection|bypass.filter|bypass.security)'
    r'|\b(act\s+as\s+a\s+|you\s+are\s+now\s+a|pretend\s+(you\s+are|to\s+be))'
    r'|\b(override|reset)\s+your\s+(system|instruction|role)',
    re.IGNORECASE,
)


def _redact_pii(text: str):
    found = False
    for pattern, placeholder in _PII_PATTERNS:
        new_text = pattern.sub(placeholder, text)
        if new_text != text:
            found = True
            text  = new_text
    return text, found


def _is_sql_injection(text: str) -> bool:
    return bool(_SQL_RE.search(text))


def _is_guardrail_violation(text: str) -> bool:
    return bool(_GUARDRAIL_RE.search(text))


def _audit_log(entry: dict):
    try:
        with open(AUDIT_DIR / "conversation_audit.jsonl", "a", encoding="utf-8") as fh:
            fh.write(json.dumps(entry, ensure_ascii=False) + "\n")
    except Exception:
        pass


# ── Pydantic models ───────────────────────────────────────────────────────────

class RiskRequest(BaseModel):
    # Required: the user's risk question
    analysis_request:  str       = ""
    # Data sources — at least one must be provided
    database_uri:      str       = ""   # any SQLAlchemy-compatible URI
    uploaded_file_ids: List[str] = []   # IDs returned by /api/upload
    # Pipeline option
    use_langgraph:     bool      = True
    # Risk context (optional — helps specialise the analysis)
    risk_domain:       str       = ""   # e.g. "credit", "market", "liquidity", "fraud"
    entity_name:       str       = ""   # e.g. company or portfolio name


class ChatRequest(BaseModel):
    message: str


class ApproveRequest(BaseModel):
    abort: bool = False


# ── Stdout capture ────────────────────────────────────────────────────────────

class _TeeWriter:
    _NOISE = (
        "You ONLY have access to the following tools",
        '"additionalProperties"',
        "Tool Name:", "Tool Arguments:", "Tool Description:",
        '"type": "object"', '"required":', '"properties":',
        "IMPORTANT: Use the following format",
        "Use the following format in your response",
        "Thought: you should always think about what to do",
        "Action Input: the input to the action, always a single string",
        "Once all necessary information is gathered",
    )
    _BOX_RE = re.compile(r'[╭╮╰╯│─]+')

    def __init__(self, task_id: str, original):
        self._tid        = task_id
        self._original   = original
        self._skip_block = False

    def write(self, text: str):
        self._original.write(text)
        self._original.flush()

        if "You ONLY have access to the following tools" in text:
            self._skip_block = True

        if text.count('|') >= 2:
            return

        stripped = text.rstrip()
        if not stripped:
            return
        first = stripped.lstrip()

        if first.startswith(('│', '╰')):
            return

        if first.startswith('╭'):
            clean_hdr = self._BOX_RE.sub('', stripped).strip()
            if not clean_hdr:
                return
            self._emit(clean_hdr)
            return

        if self._skip_block:
            if any(first.startswith(p) for p in (
                "Thought:", "Final Answer:", "Action:", "[Watchdog]",
                "━", "=", "Phase", "Error", "WARNING",
            )):
                self._skip_block = False
            else:
                return

        if any(p in text for p in self._NOISE):
            return

        if first in ('{', '}', '[', ']'):
            return
        if re.match(r'^"[a-zA-Z_]+":\s', first):
            return

        self._emit(stripped)

    def _emit(self, message: str):
        entry = {"type": "log", "message": message, "timestamp": datetime.now().isoformat()}
        tasks[self._tid]["logs"].append(entry)
        try:
            log_queues[self._tid].put_nowait(entry)
        except Exception:
            pass

    def flush(self):
        self._original.flush()

    def isatty(self):
        return False


def _push(task_id: str, msg: str, kind: str = "log"):
    entry = {"type": kind, "message": msg, "timestamp": datetime.now().isoformat()}
    tasks[task_id]["logs"].append(entry)
    try:
        log_queues[task_id].put_nowait(entry)
    except Exception:
        pass


def _strip_md_fences(text: str) -> str:
    text = re.sub(r'^```(?:markdown|md)?\s*\n', '', text.strip())
    text = re.sub(r'\n```\s*$', '', text.strip())
    text = re.sub(r'\n```(?:markdown|md)?\n[\s\S]*?```(?:\n|$)', '\n', text)
    return text.strip()


def _list_files(directory: str, extensions: list) -> list:
    p = BASE_DIR / directory
    if not p.exists():
        return []
    items = []
    for ext in extensions:
        for f in sorted(p.glob(f"*{ext}"), key=os.path.getmtime, reverse=True):
            rel = f.relative_to(BASE_DIR).as_posix()
            items.append({
                "name":     f.name,
                "url":      f"/{rel}",
                "size":     f.stat().st_size,
                "ext":      ext.lstrip("."),
                "modified": datetime.fromtimestamp(f.stat().st_mtime).isoformat(),
            })
    return items


# ── Document text extraction ──────────────────────────────────────────────────

def _extract_pdf_text(path: str) -> str:
    try:
        import pdfplumber
        parts = []
        with pdfplumber.open(path) as pdf:
            for page in pdf.pages:
                t = page.extract_text()
                if t:
                    parts.append(t)
        return "\n\n".join(parts)
    except Exception:
        pass
    try:
        import pypdf
        r = pypdf.PdfReader(path)
        return "\n\n".join(p.extract_text() or "" for p in r.pages)
    except Exception:
        return ""


def _extract_word_text(path: str) -> str:
    try:
        from docx import Document
        doc = Document(path)
        return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except Exception:
        return ""


def _extract_pptx_text(path: str) -> str:
    try:
        from pptx import Presentation
        prs   = Presentation(path)
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            texts = [sh.text.strip() for sh in slide.shapes
                     if hasattr(sh, "text") and sh.text.strip()]
            if texts:
                parts.append(f"[Slide {i}]\n" + "\n".join(texts))
        return "\n\n".join(parts)
    except Exception:
        return ""


# ── DB merge helper ───────────────────────────────────────────────────────────

def _merge_db_sources(sqlite_uri: str, db_uri: str, task_id: str) -> str:
    """Load tables from an external DB into the SQLite file DB."""
    try:
        from sqlalchemy import create_engine, inspect as sa_inspect
        src_engine  = create_engine(db_uri, connect_args={"connect_timeout": 10})
        dest_engine = create_engine(sqlite_uri)
        inspector   = sa_inspect(src_engine)
        tables      = inspector.get_table_names()[:10]
        for tbl in tables:
            try:
                df       = pd.read_sql(f"SELECT * FROM {tbl} LIMIT 50000", src_engine)
                dest_name = "db_" + "".join(c if c.isalnum() else "_" for c in tbl)
                df.to_sql(dest_name, dest_engine, if_exists="replace", index=False)
                _push(task_id, f"  Merged DB table '{tbl}' → '{dest_name}' ({len(df)} rows)", "info")
            except Exception as e:
                _push(task_id, f"  Could not merge table '{tbl}': {e}", "warn")
        return sqlite_uri
    except Exception as e:
        _push(task_id, f"  DB merge failed ({e}) — using file data only", "warn")
        return sqlite_uri


# ── Main pipeline ─────────────────────────────────────────────────────────────

def _run_analysis(task_id: str, req: RiskRequest):
    original_stdout = sys.stdout
    sys.stdout = _TeeWriter(task_id, original_stdout)

    try:
        _push(task_id, "━" * 62, "info")
        _push(task_id, "  AI FINANCIAL RISK & COMPLIANCE MONITOR  —  starting", "info")
        _push(task_id, "━" * 62, "info")

        # Clear session chart metadata
        _meta = BASE_DIR / "outputs" / "charts" / "session_charts.json"
        try:
            if _meta.exists():
                _meta.unlink()
        except Exception:
            pass

        # ── Resolve data sources ───────────────────────────────────────────────
        _structured_files = []
        _document_files   = []

        for fid in (req.uploaded_file_ids or []):
            info = _uploaded_files.get(fid.strip())
            if not info:
                continue
            if info["file_type"] == "structured":
                _structured_files.append(info)
            elif info["file_type"] == "document":
                _document_files.append(info)

        _has_explicit_db = bool(req.database_uri.strip())
        _has_structured  = bool(_structured_files)
        _has_documents   = bool(_document_files)

        if not _has_explicit_db and not _has_structured and not _has_documents:
            raise ValueError(
                "No data source provided. Please supply a Database URI or upload a file."
            )

        # ── Build unified SQLite when structured files exist ───────────────────
        _active_db: str | None = None
        if _has_structured or _has_explicit_db:
            if _has_structured:
                from sqlalchemy import create_engine as _ce
                ts_str   = datetime.now().strftime("%Y%m%d_%H%M%S")
                sq_path  = BASE_DIR / "outputs" / "uploads" / f"risk_{ts_str}.db"
                dest_eng = _ce(f"sqlite:///{sq_path}")
                for sfi in _structured_files:
                    try:
                        src_eng = _ce(sfi["sqlite_uri"])
                        df_tmp  = pd.read_sql(f"SELECT * FROM {sfi['table_name']}", src_eng)
                        df_tmp.to_sql(sfi["table_name"], dest_eng, if_exists="replace", index=False)
                        _push(task_id,
                              f"  Loaded '{sfi['table_name']}' from {sfi['original_name']} "
                              f"({len(df_tmp)} rows)", "info")
                    except Exception as _e:
                        _push(task_id, f"  Skipped '{sfi['original_name']}': {_e}", "warn")
                _active_db = f"sqlite:///{sq_path}"
                if _has_explicit_db:
                    _active_db = _merge_db_sources(_active_db, req.database_uri.strip(), task_id)
            else:
                _active_db = req.database_uri.strip()

        if _active_db:
            os.environ["ACTIVE_DB_URI"] = _active_db
        elif "ACTIVE_DB_URI" in os.environ:
            del os.environ["ACTIVE_DB_URI"]

        # ── Build document context ─────────────────────────────────────────────
        _doc_context = ""
        if _has_documents:
            parts = []
            for di in _document_files:
                if di.get("extracted_text"):
                    parts.append(
                        f"### {di['original_name']} ({di['format'].upper()})\n"
                        f"{di['extracted_text']}"
                    )
            _doc_context = "\n\n---\n\n".join(parts)

        # ── Build structured-file summary ──────────────────────────────────────
        _struct_summary = ""
        if _has_structured:
            _struct_summary = "\n".join(
                f"  • {s['original_name']} → table '{s['table_name']}' "
                f"({s['rows']} rows, {len(s['columns'])} cols: {', '.join(s['columns'][:15])})"
                for s in _structured_files
            )
        if _has_explicit_db:
            _struct_summary += f"\n  • Database URI: {req.database_uri.strip()[:80]}…"

        # ── Log data source summary ────────────────────────────────────────────
        src_parts = []
        if _has_explicit_db:  src_parts.append("Database URI")
        if _has_structured:   src_parts.append(f"{len(_structured_files)} structured file(s)")
        if _has_documents:    src_parts.append(f"{len(_document_files)} document(s)")
        _push(task_id, f"  Data sources: {' + '.join(src_parts)}", "info")

        # ── Build analysis request ─────────────────────────────────────────────
        analysis_request = req.analysis_request.strip()
        if _struct_summary:
            analysis_request += f"\n\n[AVAILABLE STRUCTURED DATA]:\n{_struct_summary}"
        if _doc_context:
            analysis_request += (
                f"\n\n[DOCUMENT CONTEXT (extracted text)]:\n{_doc_context[:8000]}"
            )
        if req.risk_domain:
            analysis_request += f"\n\n[RISK DOMAIN]: {req.risk_domain}"
        if req.entity_name:
            analysis_request += f"\n[ENTITY]: {req.entity_name}"

        # ══════════════════════════════════════════════════════════════════════
        # PHASE 1 — LangGraph Risk Planning
        # ══════════════════════════════════════════════════════════════════════
        langgraph_plan = None
        if req.use_langgraph:
            print("\n" + "=" * 62)
            print("  PHASE 1 — LangGraph Risk Planning (5 nodes)")
            print("=" * 62)

            try:
                from graphs.risk_graph import run_risk_planning
                graph_state = run_risk_planning(analysis_request, req.risk_domain)

                langgraph_plan = {
                    "analysis_plan":       graph_state.get("analysis_plan", ""),
                    "etl_guidance":        graph_state.get("etl_guidance", ""),
                    "risk_modelling":      graph_state.get("risk_modelling", ""),
                    "risk_analysis":       graph_state.get("risk_analysis", ""),
                    "preliminary_report":  graph_state.get("report_content", ""),
                }
                tasks[task_id]["langgraph_plan"] = langgraph_plan
                print("\n  LangGraph risk planning complete.\n")
            except Exception as _lg_err:
                import traceback as _tb
                _push(task_id, f"  ⚠ Phase 1 LangGraph error (non-fatal): {_lg_err}", "warn")
                _push(task_id, _tb.format_exc(), "warn")
                _push(task_id, "  Skipping Phase 1 — proceeding directly to Phase 2.", "warn")
                langgraph_plan = None

            # ── HITL Pause (only if Phase 1 succeeded) ────────────────────────
            if langgraph_plan is not None:
                _push(task_id, "━" * 58, "info")
                _push(task_id, "  HUMAN IN THE LOOP — Review the risk assessment plan", "info")
                _push(task_id, "  Approve in the UI to continue to Phase 2.", "info")
                _push(task_id, "  Auto-approving in 10 minutes if no response.", "info")
                _push(task_id, "━" * 58, "info")

                hitl_entry = {
                    "type":               "hitl_pause",
                    "task_id":            task_id,
                    "plan_preview":       langgraph_plan.get("analysis_plan", ""),
                    "preliminary_report": langgraph_plan.get("preliminary_report", ""),
                    "timestamp":          datetime.now().isoformat(),
                    "auto_approve_secs":  600,
                }
                tasks[task_id]["logs"].append(hitl_entry)
                try:
                    log_queues[task_id].put_nowait(hitl_entry)
                except Exception:
                    pass

                tasks[task_id]["status"] = "awaiting_approval"
                approved = hitl_events[task_id].wait(timeout=600)

                if tasks[task_id].get("hitl_aborted", False):
                    raise RuntimeError("Analysis aborted by user during HITL review.")

                if not approved:
                    _push(task_id, "  HITL: 10-minute timeout — auto-approving.", "warn")
                else:
                    _push(task_id, "  HITL: Approved — proceeding to Phase 2.", "success")

                tasks[task_id]["status"] = "running"

        # ══════════════════════════════════════════════════════════════════════
        # PHASE 2 — Risk Analytics CrewAI
        # ══════════════════════════════════════════════════════════════════════
        print("\n" + "=" * 62)
        print("  PHASE 2 — Risk Analytics CrewAI (5 agents)")
        print("=" * 62)

        from crew_risk import run_risk_crew
        risk_result      = run_risk_crew(analysis_request, langgraph_plan)
        risk_crew_output = risk_result.get("crew_output", "")
        analyst_output   = risk_result.get("analyst_output", "")
        scientist_output = risk_result.get("scientist_output", "")

        # Load chart metadata
        risk_blog_sections: list = []
        if _meta.exists():
            try:
                risk_blog_sections = json.loads(_meta.read_text(encoding="utf-8"))
                _seen: set = set()
                _deduped:  list = []
                for _s in risk_blog_sections:
                    _url = _s.get("chart_url", "")
                    if _url not in _seen:
                        _seen.add(_url)
                        _deduped.append(_s)
                risk_blog_sections = _deduped
                _push(task_id, f"  Risk charts loaded: {len(risk_blog_sections)}", "info")
            except Exception as _e:
                _push(task_id, f"  Chart metadata read failed: {_e}", "warn")

        risk_charts  = _list_files("outputs/charts",  [".png"])
        risk_reports = _list_files("outputs/reports", [".pdf", ".pptx", ".md"])

        risk_report = risk_crew_output if len(risk_crew_output) > 50 else \
            (langgraph_plan.get("preliminary_report", "") if langgraph_plan else
             "# Risk Analysis Complete")

        # ── Ollama enrichment ──────────────────────────────────────────────────
        risk_enhanced = risk_report
        try:
            from langchain_community.chat_models import ChatOllama
            from langchain_core.messages import HumanMessage

            if risk_report and len(risk_report) > 100:
                chart_analyses = "\n\n".join(
                    f"**{s['title']}**:\n{s['gemini_analysis']}"
                    for s in risk_blog_sections
                    if s.get("gemini_analysis")
                )
                enhance_prompt = (
                    f"You are a Chief Risk Officer preparing a board-level risk briefing.\n\n"
                    "## Draft Risk Report:\n" + risk_report + "\n\n"
                    + (f"## Chart Evidence:\n{chart_analyses}\n\n" if chart_analyses else "")
                    + "## Your Task — Targeted Enrichment:\n"
                    "Produce the final board-ready version (700–1000 words). "
                    "Keep Executive Summary max 4 bullets. "
                    "Keep Key Findings max 5 points with specific metrics. "
                    "Enrich Recommendations with Owner|Metric Target|30-day|60-day|90-day table. "
                    "Include a Regulatory Compliance Status section (Basel III / IFRS 9 / FRTB). "
                    "Replace Next Steps with numbered action plan [N]. Action — Owner — Deadline — Metric. "
                    "IMPORTANT: Do NOT include chart descriptions in the report body — "
                    "those are rendered separately in the Chart Analysis section. "
                    "Output the complete report only — no preamble."
                )
                _push(task_id, "  Ollama: enriching risk report…", "info")
                _ollama  = ChatOllama(model="qwen3.5:cloud")
                response = _ollama.invoke([HumanMessage(content=enhance_prompt)])
                enhanced = _strip_md_fences(response.content)
                if enhanced and len(enhanced) > 200:
                    risk_enhanced = enhanced
                    _push(task_id, "  Ollama: risk report enriched ✓", "success")
        except Exception as _e:
            _push(task_id, f"  Ollama enrichment (non-critical): {_e}", "warn")

        # ── Build compliance summary ────────────────────────────────────────────
        compliance_summary = ""
        try:
            from langchain_community.chat_models import ChatOllama
            from langchain_core.messages import HumanMessage as _HM

            _comp_prompt = (
                "You are a bank regulatory compliance officer.\n\n"
                f"## Risk Analysis Report:\n{risk_enhanced[:3000]}\n\n"
                "## Chart Risk Evidence:\n"
                + "\n\n".join(
                    f"**{s['title']}**: {s.get('gemini_analysis','')[:300]}"
                    for s in risk_blog_sections if s.get("gemini_analysis")
                )[:1500]
                + "\n\nWrite a structured Regulatory Compliance Summary covering:\n"
                "1. **Basel III Status** — Capital Adequacy, PD/LGD/EAD, RWA assessment\n"
                "2. **IFRS 9 Status** — ECL stage classification (Stage 1/2/3), provisioning\n"
                "3. **FRTB / Market Risk** — VaR, SVaR, ES, IMA vs SA (if applicable)\n"
                "4. **AML / Fraud Controls** — transaction monitoring, suspicious indicators\n"
                "5. **Overall Compliance Rating**: COMPLIANT / REQUIRES ACTION / CRITICAL\n"
                "6. **Immediate Actions Required** (numbered list)\n\n"
                "Rate each framework: ✅ COMPLIANT | ⚠️ REQUIRES ACTION | ❌ CRITICAL | N/A\n"
                "Output markdown. Keep to 400–600 words."
            )
            _push(task_id, "  Ollama: generating compliance summary…", "info")
            _ollama2 = ChatOllama(model="qwen3.5:cloud")
            _resp2   = _ollama2.invoke([_HM(content=_comp_prompt)])
            compliance_summary = _strip_md_fences(_resp2.content)
            if compliance_summary and len(compliance_summary) > 100:
                _push(task_id, "  Ollama: compliance summary generated ✓", "success")
        except Exception as _e:
            _push(task_id, f"  Compliance summary (non-critical): {_e}", "warn")

        tasks[task_id].update({
            "risk_report":         risk_enhanced,
            "compliance_summary":  compliance_summary,
            "risk_blog_sections":  risk_blog_sections,
            "risk_charts":         risk_charts,
            "analyst_output":      analyst_output,
            "scientist_output":    scientist_output,
            "status":              "completed",
            "completed_at":        datetime.now().isoformat(),
            "all_reports":         risk_reports,
        })

        _push(task_id, "━" * 62, "success")
        _push(task_id, f"  Analysis complete!  Charts: {len(risk_charts)}  Reports: {len(risk_reports)}", "success")
        _push(task_id, "━" * 62, "success")

    except Exception as exc:
        import traceback
        err_trace = traceback.format_exc()
        tasks[task_id].update({
            "status":       "error",
            "error":        str(exc),
            "completed_at": datetime.now().isoformat(),
        })
        _push(task_id, f"ERROR: {exc}", "error")
        _push(task_id, err_trace, "error")

    finally:
        sys.stdout = original_stdout
        done_type  = "complete" if tasks[task_id]["status"] == "completed" else (
            "aborted" if tasks[task_id].get("hitl_aborted") else "error"
        )
        try:
            log_queues[task_id].put_nowait({
                "type": done_type, "task_id": task_id,
                "timestamp": datetime.now().isoformat()
            })
        except Exception:
            pass
        hitl_events.pop(task_id, None)


def _run_with_lock(task_id: str, req: RiskRequest):
    with _analysis_lock:
        _run_analysis(task_id, req)


# ── Routes ────────────────────────────────────────────────────────────────────

@app.get("/", include_in_schema=False)
async def root():
    return FileResponse(str(BASE_DIR / "static" / "index.html"))


@app.post("/api/analyze")
async def start_analysis(req: RiskRequest):
    """Start a new financial risk analysis (background thread)."""
    if not req.analysis_request.strip():
        raise HTTPException(status_code=400,
                            detail="analysis_request is required.")
    if not req.database_uri.strip() and not req.uploaded_file_ids:
        raise HTTPException(status_code=400,
                            detail="Provide a Database URI or upload at least one file.")
    if _analysis_lock.locked():
        raise HTTPException(status_code=409,
                            detail="An analysis is already running.")

    task_id = str(uuid.uuid4())
    tasks[task_id] = {
        "status":              "running",
        "started_at":          datetime.now().isoformat(),
        "completed_at":        None,
        "entity_name":         req.entity_name,
        "risk_domain":         req.risk_domain,
        "database_uri":        req.database_uri,
        "error":               None,
        "logs":                [],
        "langgraph_plan":      {},
        "hitl_aborted":        False,
        # Risk outputs
        "risk_report":         "",
        "compliance_summary":  "",
        "risk_blog_sections":  [],
        "risk_charts":         [],
        "analyst_output":      "",
        "scientist_output":    "",
        "all_reports":         [],
        # Conversation
        "conversation":        [],
    }
    log_queues[task_id]  = queue.Queue()
    hitl_events[task_id] = threading.Event()

    thread = threading.Thread(
        target=_run_with_lock, args=(task_id, req), daemon=True
    )
    thread.start()

    return {"task_id": task_id, "status": "started"}


@app.get("/api/status/{task_id}")
async def get_status(task_id: str):
    if task_id not in tasks:
        raise HTTPException(status_code=404, detail="Task not found.")
    t = tasks[task_id]
    return {
        "task_id":      task_id,
        "status":       t["status"],
        "started_at":   t["started_at"],
        "completed_at": t.get("completed_at"),
        "error":        t.get("error"),
        "log_lines":    len(t.get("logs", [])),
    }


@app.get("/api/results/{task_id}")
async def get_results(task_id: str):
    if task_id not in tasks:
        raise HTTPException(status_code=404, detail="Task not found.")
    t = tasks[task_id]
    if t["status"] == "running":
        raise HTTPException(status_code=202, detail="Analysis still running.")
    if t["status"] == "error":
        raise HTTPException(status_code=500, detail=t.get("error", "Unknown error."))

    return {
        "task_id":            task_id,
        "status":             "completed",
        "risk_report":        t.get("risk_report", ""),
        "compliance_summary": t.get("compliance_summary", ""),
        "risk_blog_sections": t.get("risk_blog_sections", []),
        "risk_charts":        t.get("risk_charts", []),
        "all_reports":        t.get("all_reports", []),
        "langgraph_plan":     t.get("langgraph_plan", {}),
        "completed_at":       t.get("completed_at"),
    }


@app.get("/api/logs/{task_id}")
async def get_logs(task_id: str):
    if task_id not in tasks:
        raise HTTPException(status_code=404, detail="Task not found.")
    return {"task_id": task_id, "logs": tasks[task_id].get("logs", [])}


@app.get("/api/stream/{task_id}")
async def stream_events(task_id: str):
    if task_id not in tasks:
        raise HTTPException(status_code=404, detail="Task not found.")

    async def generator():
        q        = log_queues[task_id]
        buffered = list(tasks[task_id].get("logs", []))
        replayed = len(buffered)
        for entry in buffered:
            yield f"data: {json.dumps(entry)}\n\n"

        drained = 0
        while drained < replayed:
            try:
                q.get_nowait()
                drained += 1
            except queue.Empty:
                break

        while True:
            try:
                msg = q.get_nowait()
                yield f"data: {json.dumps(msg)}\n\n"
                if msg.get("type") in ("complete", "error", "aborted"):
                    return
            except queue.Empty:
                if tasks[task_id]["status"] in ("completed", "error", "aborted"):
                    final = "complete" if tasks[task_id]["status"] == "completed" else "error"
                    yield f"data: {json.dumps({'type': final, 'task_id': task_id})}\n\n"
                    return
                await asyncio.sleep(0.35)

    return StreamingResponse(
        generator(), media_type="text/event-stream",
        headers={"Cache-Control": "no-cache",
                 "X-Accel-Buffering": "no",
                 "Connection": "keep-alive"}
    )


@app.get("/api/charts")
async def list_charts():
    return {"charts": _list_files("outputs/charts", [".png"])}


@app.get("/api/reports")
async def list_reports():
    return {"reports": _list_files("outputs/reports", [".pdf", ".pptx", ".md"])}


@app.get("/api/tasks")
async def list_tasks():
    return [
        {"task_id": tid, "status": t["status"],
         "started_at": t["started_at"], "completed_at": t.get("completed_at"),
         "entity_name": t.get("entity_name", ""), "risk_domain": t.get("risk_domain", "")}
        for tid, t in tasks.items()
    ]


@app.post("/api/approve/{task_id}")
async def approve_hitl(task_id: str, req: ApproveRequest = None):
    if req is None:
        req = ApproveRequest()
    if task_id not in tasks:
        raise HTTPException(status_code=404, detail="Task not found.")
    if tasks[task_id]["status"] != "awaiting_approval":
        raise HTTPException(status_code=400, detail="Task is not awaiting approval.")
    if task_id not in hitl_events:
        raise HTTPException(status_code=400, detail="No HITL event found.")
    if req.abort:
        tasks[task_id]["hitl_aborted"] = True
    hitl_events[task_id].set()
    return {"task_id": task_id, "aborted": req.abort}


# ── File upload ───────────────────────────────────────────────────────────────

_SUPPORTED_EXTS = {
    ".csv", ".xlsx", ".xls",          # structured data
    ".pdf", ".docx", ".doc",           # documents
    ".pptx", ".ppt",                   # presentations
}
_STRUCTURED_EXTS = {".csv", ".xlsx", ".xls"}


@app.post("/api/upload")
async def upload_data_files(files: List[UploadFile] = File(...)):
    """
    Upload one or more files for risk analysis.
    CSV/Excel → loaded into SQLite (structured data).
    PDF/Word/PPTX → text extracted as document context.
    """
    if not files:
        raise HTTPException(status_code=400, detail="No files provided.")

    upload_dir = BASE_DIR / "outputs" / "uploads"
    upload_dir.mkdir(parents=True, exist_ok=True)

    results = []
    for file in files:
        if not file.filename:
            continue
        ext = Path(file.filename).suffix.lower()
        if ext not in _SUPPORTED_EXTS:
            raise HTTPException(
                status_code=400,
                detail=(f"'{file.filename}': unsupported format. "
                        "Allowed: CSV, Excel, PDF, Word (.docx), PowerPoint (.pptx).")
            )

        ts        = datetime.now().strftime("%Y%m%d_%H%M%S_%f")
        safe_name = "".join(c if c.isalnum() or c in "._-" else "_" for c in file.filename)
        raw_path  = upload_dir / f"{ts}_{safe_name}"
        content   = await file.read()
        raw_path.write_bytes(content)

        file_id = ts

        if ext in _STRUCTURED_EXTS:
            try:
                df = pd.read_csv(raw_path) if ext == ".csv" else pd.read_excel(raw_path)
            except Exception as e:
                raise HTTPException(status_code=422,
                                    detail=f"Could not read '{file.filename}': {e}")

            tbl = Path(file.filename).stem.lower()
            tbl = ("".join(c if c.isalnum() else "_" for c in tbl))[:50] or "data"

            sqlite_path = upload_dir / f"{ts}_data.db"
            try:
                from sqlalchemy import create_engine as _ce
                pd.DataFrame(df).to_sql(
                    tbl, _ce(f"sqlite:///{sqlite_path}"), if_exists="replace", index=False
                )
            except Exception as e:
                raise HTTPException(status_code=500,
                                    detail=f"Could not store '{file.filename}': {e}")

            _uploaded_files[file_id] = {
                "original_name":  file.filename,
                "file_type":      "structured",
                "format":         ext.lstrip("."),
                "table_name":     tbl,
                "sqlite_uri":     f"sqlite:///{sqlite_path}",
                "extracted_text": None,
                "rows":           len(df),
                "columns":        list(df.columns),
            }
            results.append({
                "file_id":    file_id,
                "filename":   file.filename,
                "type":       "structured",
                "format":     ext.lstrip("."),
                "table_name": tbl,
                "rows":       len(df),
                "columns":    list(df.columns),
            })

        else:
            if ext == ".pdf":
                text = _extract_pdf_text(str(raw_path))
            elif ext in (".docx", ".doc"):
                text = _extract_word_text(str(raw_path))
            else:
                text = _extract_pptx_text(str(raw_path))

            word_count = len(text.split()) if text else 0
            _uploaded_files[file_id] = {
                "original_name":  file.filename,
                "file_type":      "document",
                "format":         ext.lstrip("."),
                "table_name":     None,
                "sqlite_uri":     None,
                "extracted_text": text[:60000],
                "rows":           None,
                "columns":        None,
            }
            results.append({
                "file_id":    file_id,
                "filename":   file.filename,
                "type":       "document",
                "format":     ext.lstrip("."),
                "word_count": word_count,
                "preview":    text[:200] if text else "",
            })

    return {"files": results, "count": len(results)}


# ── Agentic Chat ──────────────────────────────────────────────────────────────

_VIZ_KEYWORDS = (
    "visualiz", "chart", "graph", "plot", "bar chart", "pie chart",
    "histogram", "scatter", "heatmap", "heat map", "line chart",
    "show me a", "draw", "figure",
)


async def _agentic_chat(message: str, task_ctx: dict,
                        history: list, task_id: str) -> tuple:
    """NVIDIA LLaMA agentic chat with intent-based visualization."""
    from config import query_nvidia, nvidia_client
    if not nvidia_client:
        return "NVIDIA API key not configured.", [], []

    risk_report       = (task_ctx.get("risk_report") or "")[:4000]
    compliance_summary = (task_ctx.get("compliance_summary") or "")[:1500]
    chart_analyses    = "\n\n".join(
        f"**{s['title']}**:\n{s.get('gemini_analysis', '')}"
        for s in (task_ctx.get("risk_blog_sections") or [])
        if s.get("gemini_analysis")
    )[:2000]

    context_block = (
        "## Risk Assessment Report:\n" + risk_report + "\n\n"
        + (f"## Chart Analyses:\n{chart_analyses}\n\n" if chart_analyses else "")
        + "## Regulatory Compliance Summary:\n" + compliance_summary
    )

    msg_lower = message.lower()
    wants_viz = any(kw in msg_lower for kw in _VIZ_KEYWORDS)
    chart_urls: list = []

    if wants_viz:
        code_system = (
            "You are a Python data visualization expert.\n"
            "Output a single ```python ... ``` code block and nothing else.\n\n"
            "STRICT RULES:\n"
            "1. Use ONLY: matplotlib, seaborn, pandas, numpy, json, datetime\n"
            "2. Do NOT import os, sys, subprocess, requests, pathlib or any other module\n"
            "3. Save with exactly: plt.savefig(OUTPUT_PATH, dpi=150, bbox_inches='tight')\n"
            "4. Do NOT call plt.show()\n"
            "5. Use hardcoded/synthetic data matching the risk report context\n"
            "6. Chart must be dark-themed (plt.style.use('dark_background') already applied)\n"
            "7. Add a clear title and axis labels"
        )
        code_user = (
            f"{context_block}\n\nUser request: {message}\n\n"
            "Write the Python visualization code now. Output ONLY the ```python``` block."
        )
        try:
            code_response = await asyncio.to_thread(
                query_nvidia,
                [{"role": "system", "content": code_system},
                 {"role": "user",   "content": code_user}],
                0.1, 1800
            )
        except Exception as exc:
            return f"NVIDIA error generating visualization: {exc}", [], []

        code_match = re.search(r"```python\s*([\s\S]*?)```", code_response, re.IGNORECASE)
        if not code_match:
            code_match = re.search(r"(import matplotlib[\s\S]*)", code_response)

        if code_match:
            code           = code_match.group(1).strip()
            ts_viz         = datetime.now().strftime("%Y%m%d_%H%M%S_%f")
            chart_filename = f"chat_viz_{task_id[:8]}_{ts_viz}.png"
            chart_path     = str(BASE_DIR / "outputs" / "charts" / chart_filename)
            chart_url      = f"/outputs/charts/{chart_filename}"

            # Safe code execution
            BLOCKED = ["import os", "import sys", "import subprocess", "import socket",
                       "import requests", "import pathlib", "__import__", "eval(", "exec(",
                       "open(", "os.system", "os.popen", "shutil"]
            blocked_found = next((b for b in BLOCKED if b in code), None)
            if blocked_found:
                viz_status = f"Blocked pattern '{blocked_found}' not allowed."
            else:
                setup = (
                    "import matplotlib\nmatplotlib.use('Agg')\n"
                    "import matplotlib.pyplot as plt\nimport pandas as pd\nimport numpy as np\n"
                    "import json\nimport seaborn as sns\nfrom datetime import datetime\n"
                    f"OUTPUT_PATH = {repr(chart_path)}\n"
                    "plt.style.use('dark_background')\n"
                    "plt.rcParams.update({"
                    "'figure.facecolor':'#0d1117','axes.facecolor':'#161b22',"
                    "'text.color':'#e6edf3','axes.labelcolor':'#e6edf3',"
                    "'xtick.color':'#8b949e','ytick.color':'#8b949e',"
                    "'axes.edgecolor':'#30363d','grid.color':'#21262d'"
                    "})\n"
                )
                full_code = setup + "\n" + code
                if "plt.show()" in full_code:
                    full_code = full_code.replace(
                        "plt.show()",
                        'plt.savefig(OUTPUT_PATH, dpi=120, bbox_inches="tight"); plt.close()'
                    )
                if "plt.savefig" not in full_code:
                    full_code += '\nplt.savefig(OUTPUT_PATH, dpi=120, bbox_inches="tight")\nplt.close()\n'
                try:
                    exec(full_code, {})  # noqa: S102
                    if Path(chart_path).exists():
                        chart_urls.append(chart_url)
                        viz_status = f"Chart generated at {chart_url}."
                    else:
                        viz_status = "Chart file was not created."
                except Exception as exc:
                    viz_status = f"Code execution error: {exc}"
        else:
            viz_status = "Could not extract Python code from model response."

        explain_messages = [
            {"role": "system", "content": (
                "You are an expert risk analyst assistant. "
                "Answer based on the risk analytics context. "
                "Be concise and insightful. Do NOT generate code."
            )}
        ]
        for turn in history[-4:]:
            explain_messages.append({"role": "user",      "content": turn["user"]})
            explain_messages.append({"role": "assistant", "content": turn["assistant"]})
        explain_messages.append({"role": "user", "content": (
            f"{context_block}\n\nUser asked: {message}\n\n"
            f"Visualization status: {viz_status}\n\n"
            "Provide a concise explanation of what the chart shows and key risk insights."
        )})
        try:
            answer = await asyncio.to_thread(query_nvidia, explain_messages, 0.3, 1024)
        except Exception:
            answer = f"Chart generated. {viz_status}"

        return answer, chart_urls, []

    # Conversational answer
    conv_messages = [{"role": "system", "content": (
        "You are an expert financial risk analyst AI assistant.\n"
        "Answer questions strictly based on the risk analytics results provided.\n"
        "Be concise, specific, and reference actual data from the report.\n"
        "Do NOT generate code or charts — just answer directly.\n\n"
        + context_block
    )}]
    for turn in history[-6:]:
        conv_messages.append({"role": "user",      "content": turn["user"]})
        conv_messages.append({"role": "assistant", "content": turn["assistant"]})
    conv_messages.append({"role": "user", "content": message})

    try:
        answer = await asyncio.to_thread(query_nvidia, conv_messages, 0.3, 1500)
    except Exception as exc:
        return f"NVIDIA LLaMA error: {exc}", [], []

    return answer, [], []


@app.post("/api/chat/{task_id}")
async def chat(task_id: str, req: ChatRequest):
    """Agentic chat powered by NVIDIA LLaMA with on-demand visualization."""
    if task_id not in tasks:
        raise HTTPException(status_code=404, detail="Task not found.")
    t = tasks[task_id]
    if t["status"] != "completed":
        raise HTTPException(status_code=400, detail="Analysis not yet completed.")

    message = req.message.strip()
    if len(message) > 2000:
        raise HTTPException(status_code=400, detail="Message too long (max 2000 chars).")
    if _is_sql_injection(message):
        _audit_log({"timestamp": datetime.now().isoformat(), "task_id": task_id,
                    "flags": ["sql_injection"], "blocked": True})
        raise HTTPException(status_code=400, detail="Blocked: SQL injection pattern.")
    if _is_guardrail_violation(message):
        _audit_log({"timestamp": datetime.now().isoformat(), "task_id": task_id,
                    "flags": ["guardrail_violation"], "blocked": True})
        raise HTTPException(status_code=400, detail="Blocked: policy violation.")

    clean_message, pii_in_input = _redact_pii(message)

    answer, chart_urls, _ = await _agentic_chat(
        message  = clean_message,
        task_ctx = t,
        history  = t.get("conversation", [])[-10:],
        task_id  = task_id,
    )

    reply, pii_in_output = _redact_pii(answer)
    security_flags = []
    if pii_in_input:  security_flags.append("pii_redacted_input")
    if pii_in_output: security_flags.append("pii_redacted_output")

    _audit_log({"timestamp": datetime.now().isoformat(), "task_id": task_id,
                "input_len": len(message), "flags": security_flags, "blocked": False,
                "output_len": len(reply), "chart_urls": chart_urls})

    t["conversation"].append({"user": clean_message, "assistant": reply,
                               "chart_urls": chart_urls})
    return {"reply": reply, "chart_urls": chart_urls, "security_flags": security_flags}


@app.get("/api/chat/history/{task_id}")
async def chat_history(task_id: str):
    if task_id not in tasks:
        raise HTTPException(status_code=404, detail="Task not found.")
    return {"history": tasks[task_id].get("conversation", [])}


if __name__ == "__main__":
    import uvicorn
    uvicorn.run("app:app", host="127.0.0.1", port=8003, reload=False)
