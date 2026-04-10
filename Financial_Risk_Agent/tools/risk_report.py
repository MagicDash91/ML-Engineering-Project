"""
tools/risk_report.py – Professional PDF and PowerPoint report generation
for financial risk and compliance assessments.
"""

import os
import sys
import re
import ast
import json
from datetime import datetime

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from crewai.tools import tool

BASE_DIR    = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
REPORTS_DIR = os.path.join(BASE_DIR, "outputs", "reports")
os.makedirs(REPORTS_DIR, exist_ok=True)


@tool("Generate Risk PDF Report")
def generate_risk_pdf_report(
    report_content: str,
    chart_paths: str = "[]",
    report_title: str = "Financial Risk Assessment Report",
) -> str:
    """
    Generate a professional A4 PDF with risk report text and embedded chart images.

    Args:
        report_content: Markdown-formatted risk report text.
        chart_paths:    JSON array of PNG file paths to include.
        report_title:   PDF title shown on the cover page.

    Returns:
        JSON with 'pdf_path' and 'charts_included' count.
    """
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import inch, cm
        from reportlab.lib import colors
        from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer,
                                        Image, HRFlowable, PageBreak)
        from reportlab.lib.enums import TA_CENTER, TA_JUSTIFY

        ts       = datetime.now().strftime("%Y%m%d_%H%M%S")
        pdf_path = os.path.join(REPORTS_DIR, f"risk_report_{ts}.pdf")

        doc  = SimpleDocTemplate(pdf_path, pagesize=A4,
                                  rightMargin=2*cm, leftMargin=2*cm,
                                  topMargin=2.5*cm, bottomMargin=2*cm)
        base   = getSampleStyleSheet()
        C_DARK = colors.HexColor("#0f172a")
        C_RISK = colors.HexColor("#dc2626")   # risk red
        C_BLUE = colors.HexColor("#1d4ed8")

        s_title = ParagraphStyle("s_title", parent=base["Title"], fontSize=22,
                                  textColor=C_DARK, spaceAfter=18,
                                  alignment=TA_CENTER, fontName="Helvetica-Bold")
        s_h1    = ParagraphStyle("s_h1", parent=base["Heading1"], fontSize=14,
                                  textColor=C_RISK, spaceBefore=16, spaceAfter=8,
                                  fontName="Helvetica-Bold")
        s_h2    = ParagraphStyle("s_h2", parent=base["Heading2"], fontSize=12,
                                  textColor=C_BLUE, spaceBefore=10, spaceAfter=5,
                                  fontName="Helvetica-Bold")
        s_body  = ParagraphStyle("s_body", parent=base["Normal"], fontSize=10,
                                  leading=14, spaceAfter=7,
                                  alignment=TA_JUSTIFY, fontName="Helvetica")
        s_sub   = ParagraphStyle("s_sub", parent=base["Normal"], fontSize=10,
                                  alignment=TA_CENTER, textColor=colors.grey)

        story = []
        story.append(Spacer(1, 0.6*inch))
        story.append(Paragraph(report_title, s_title))
        story.append(Paragraph(
            f"Generated: {datetime.now().strftime('%B %d, %Y  %H:%M')}  •  CONFIDENTIAL",
            s_sub))
        story.append(HRFlowable(width="100%", thickness=2, color=C_RISK))
        story.append(Spacer(1, 0.4*inch))

        _bold_re = re.compile(r"\*\*(.+?)\*\*")
        for raw_line in report_content.split("\n"):
            line = raw_line.strip()
            if not line:
                story.append(Spacer(1, 0.08*inch))
            elif line.startswith("# "):
                story.append(Paragraph(line[2:], s_title))
            elif line.startswith("## "):
                story.append(Paragraph(line[3:], s_h1))
            elif line.startswith("### "):
                story.append(Paragraph(line[4:], s_h2))
            elif line.startswith(("- ", "* ", "• ")):
                body = _bold_re.sub(r"<b>\1</b>", line[2:])
                story.append(Paragraph(f"&bull;&nbsp;{body}", s_body))
            else:
                story.append(Paragraph(_bold_re.sub(r"<b>\1</b>", line), s_body))

        chart_list: list = []
        if chart_paths:
            try:
                chart_list = json.loads(chart_paths)
            except (json.JSONDecodeError, ValueError):
                try:
                    chart_list = ast.literal_eval(chart_paths)
                except Exception:
                    pass

        if chart_list:
            story.append(PageBreak())
            story.append(Paragraph("Risk Visualisations", s_h1))
            story.append(HRFlowable(width="100%", thickness=1, color=C_BLUE))
            story.append(Spacer(1, 0.2*inch))
            for cp in chart_list:
                if os.path.exists(cp):
                    story.append(Image(cp, width=6.5*inch, height=4.0*inch))
                    story.append(Spacer(1, 0.25*inch))

        doc.build(story)

        return json.dumps({
            "status":           "success",
            "pdf_path":         pdf_path,
            "charts_included":  len([c for c in chart_list if os.path.exists(c)]),
        }, indent=2)
    except Exception as e:
        return json.dumps({"status": "error", "message": str(e)})


@tool("Generate Risk PowerPoint Presentation")
def generate_risk_ppt_report(
    report_content: str,
    chart_paths: str = "[]",
    report_title: str = "Financial Risk Assessment",
) -> str:
    """
    Generate a widescreen (16:9) PowerPoint risk presentation with branded slides.

    Args:
        report_content: Markdown text — each '## Section' becomes a content slide.
        chart_paths:    JSON array of PNG paths — each becomes a full-bleed chart slide.
        report_title:   Presentation title on the cover slide.

    Returns:
        JSON with 'ppt_path' and 'slides_created' count.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        prs = Presentation()
        prs.slide_width  = Inches(13.33)
        prs.slide_height = Inches(7.5)

        C_DARK    = RGBColor(0x0f, 0x17, 0x2a)
        C_RISK    = RGBColor(0xdc, 0x26, 0x26)
        C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
        C_TEXT    = RGBColor(0x1e, 0x29, 0x3b)
        C_SUBTEXT = RGBColor(0xFE, 0xCA, 0xCA)
        BLANK     = prs.slide_layouts[6]

        def _rect(slide, l, t, w, h, color):
            shp = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
            shp.fill.solid()
            shp.fill.fore_color.rgb = color
            shp.line.fill.background()
            return shp

        def _textbox(slide, l, t, w, h, text, size,
                     bold=False, color=C_WHITE, align=PP_ALIGN.LEFT):
            tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
            tf = tb.text_frame
            tf.word_wrap = True
            p   = tf.paragraphs[0]
            p.text      = text
            p.alignment = align
            run = p.runs[0] if p.runs else p.add_run()
            run.font.size      = Pt(size)
            run.font.bold      = bold
            run.font.color.rgb = color

        def add_cover(title, subtitle):
            slide = prs.slides.add_slide(BLANK)
            bg    = slide.background.fill
            bg.solid()
            bg.fore_color.rgb = C_DARK
            _rect(slide, 0, 3.2, 13.33, 0.06, C_RISK)
            _textbox(slide, 1.0, 1.5, 11.33, 1.8, title,
                     34, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
            _textbox(slide, 1.0, 3.5, 11.33, 0.9, subtitle,
                     16, color=C_SUBTEXT, align=PP_ALIGN.CENTER)

        def add_content(section_title, lines):
            slide = prs.slides.add_slide(BLANK)
            _rect(slide, 0, 0, 13.33, 1.15, C_DARK)
            _rect(slide, 0, 1.15, 0.07, 6.35, C_RISK)
            _textbox(slide, 0.3, 0.15, 12.73, 0.85, section_title, 22, bold=True)
            tb = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.35), Inches(12.33), Inches(5.7))
            tf = tb.text_frame
            tf.word_wrap = True
            for i, ln in enumerate(lines[:16]):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.space_before = Pt(5)
                cleaned = ln.replace("**", "").replace("##", "").replace("#", "").strip()
                p.text  = cleaned
                run     = p.runs[0] if p.runs else p.add_run()
                run.font.size      = Pt(13 if ln.startswith("**") else 12)
                run.font.bold      = ln.startswith("**")
                run.font.color.rgb = C_TEXT

        def add_chart_slide(chart_path, label):
            slide = prs.slides.add_slide(BLANK)
            _rect(slide, 0, 0, 13.33, 1.05, C_DARK)
            _rect(slide, 0, 1.05, 0.07, 6.45, C_RISK)
            _textbox(slide, 0.3, 0.12, 12.73, 0.82, label, 20, bold=True)
            if os.path.exists(chart_path):
                slide.shapes.add_picture(
                    chart_path, Inches(0.5), Inches(1.1), Inches(12.33), Inches(6.0))

        add_cover(
            report_title,
            f"Financial Risk & Compliance Assessment  •  {datetime.now().strftime('%B %Y')}  •  CONFIDENTIAL"
        )

        _normalised = re.sub(r'\n#{2,3}\s+', '\n## ', report_content)
        sections    = _normalised.split("\n## ")
        for sec in sections[:8]:
            raw_lines = sec.strip().split("\n")
            sec_title = raw_lines[0].replace("#", "").strip()
            body      = [l.strip() for l in raw_lines[1:] if l.strip()]
            if sec_title and body:
                add_content(sec_title, body)

        chart_list: list = []
        if chart_paths:
            try:
                chart_list = json.loads(chart_paths)
            except (json.JSONDecodeError, ValueError):
                try:
                    chart_list = ast.literal_eval(chart_paths)
                except Exception:
                    pass

        for i, cp in enumerate(chart_list[:6]):
            label = os.path.basename(cp).replace(".png", "").replace("_", " ").title()
            add_chart_slide(cp, f"Risk Chart {i + 1}: {label}")

        add_cover("End of Report", "Confidential — For Internal Use Only")

        ts       = datetime.now().strftime("%Y%m%d_%H%M%S")
        ppt_path = os.path.join(REPORTS_DIR, f"risk_presentation_{ts}.pptx")
        prs.save(ppt_path)

        return json.dumps({
            "status":        "success",
            "ppt_path":      ppt_path,
            "slides_created": len(prs.slides),
        }, indent=2)
    except Exception as e:
        return json.dumps({"status": "error", "message": str(e)})
