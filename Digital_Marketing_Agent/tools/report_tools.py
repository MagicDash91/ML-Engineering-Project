"""
tools/report_tools.py – PDF, PPTX, and Markdown report generation tools.
"""
from __future__ import annotations

import re
from datetime import datetime
from pathlib import Path

from crewai.tools import tool

_REPORTS = Path(__file__).parent.parent / "outputs" / "reports"
_REPORTS.mkdir(parents=True, exist_ok=True)


def _safe_filename(brand: str) -> str:
    return "".join(c if c.isalnum() else "_" for c in brand)


# ── Markdown ──────────────────────────────────────────────────────────────────

@tool("generate_text_report")
def generate_text_report(content: str, brand_name: str) -> str:
    """Save a markdown campaign report to outputs/reports/.

    Args:
        content: Markdown content for the report.
        brand_name: Brand name used in filename.

    Returns:
        Path to the saved .md file.
    """
    ts   = datetime.now().strftime("%Y%m%d_%H%M%S")
    path = _REPORTS / f"campaign_report_{_safe_filename(brand_name)}_{ts}.md"
    path.write_text(content, encoding="utf-8")
    print(f"[ReportTools] Markdown report saved → {path}", flush=True)
    return str(path)


# ── PDF via ReportLab ─────────────────────────────────────────────────────────

@tool("generate_pdf_report")
def generate_pdf_report(content: str, brand_name: str) -> str:
    """Generate a PDF campaign report using ReportLab.

    Args:
        content: Markdown/plain-text content for the PDF.
        brand_name: Brand name used in filename and title.

    Returns:
        Path to the saved .pdf file.
    """
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import cm
        from reportlab.lib import colors
        from reportlab.platypus import (
            SimpleDocTemplate, Paragraph, Spacer, HRFlowable
        )
        from reportlab.lib.enums import TA_LEFT, TA_CENTER
    except ImportError as exc:
        # Fallback: save as .txt
        ts   = datetime.now().strftime("%Y%m%d_%H%M%S")
        path = _REPORTS / f"campaign_report_{_safe_filename(brand_name)}_{ts}.txt"
        path.write_text(content, encoding="utf-8")
        return f"[ReportLab not installed – saved as text: {path}]"

    ts       = datetime.now().strftime("%Y%m%d_%H%M%S")
    pdf_path = _REPORTS / f"campaign_report_{_safe_filename(brand_name)}_{ts}.pdf"

    doc = SimpleDocTemplate(
        str(pdf_path),
        pagesize=A4,
        leftMargin=2*cm, rightMargin=2*cm,
        topMargin=2*cm, bottomMargin=2*cm,
    )

    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        "TitleCustom", parent=styles["Title"],
        fontSize=20, textColor=colors.HexColor("#1a73e8"), spaceAfter=12,
    )
    h1_style = ParagraphStyle(
        "H1Custom", parent=styles["Heading1"],
        fontSize=14, textColor=colors.HexColor("#1a73e8"), spaceBefore=12, spaceAfter=6,
    )
    h2_style = ParagraphStyle(
        "H2Custom", parent=styles["Heading2"],
        fontSize=12, textColor=colors.HexColor("#333333"), spaceBefore=8, spaceAfter=4,
    )
    body_style = ParagraphStyle(
        "BodyCustom", parent=styles["Normal"],
        fontSize=10, leading=15, spaceAfter=6,
    )

    story = []
    story.append(Paragraph(f"Campaign Report – {brand_name}", title_style))
    story.append(Paragraph(f"Generated: {datetime.now().strftime('%B %d, %Y')}", styles["Normal"]))
    story.append(HRFlowable(width="100%", thickness=1, color=colors.HexColor("#1a73e8")))
    story.append(Spacer(1, 12))

    # Parse markdown lines into ReportLab flowables
    for line in content.splitlines():
        stripped = line.strip()
        if not stripped:
            story.append(Spacer(1, 6))
        elif stripped.startswith("## "):
            story.append(Paragraph(stripped[3:], h1_style))
        elif stripped.startswith("### "):
            story.append(Paragraph(stripped[4:], h2_style))
        elif stripped.startswith("# "):
            story.append(Paragraph(stripped[2:], title_style))
        elif stripped.startswith("- ") or stripped.startswith("* "):
            story.append(Paragraph(f"• {stripped[2:]}", body_style))
        elif re.match(r"^\d+\. ", stripped):
            story.append(Paragraph(stripped, body_style))
        else:
            # Strip inline markdown
            clean = re.sub(r"\*\*(.+?)\*\*", r"<b>\1</b>", stripped)
            clean = re.sub(r"\*(.+?)\*", r"<i>\1</i>", clean)
            story.append(Paragraph(clean, body_style))

    doc.build(story)
    print(f"[ReportTools] PDF report saved → {pdf_path}", flush=True)
    return str(pdf_path)


# ── PowerPoint via python-pptx ────────────────────────────────────────────────

@tool("generate_ppt_report")
def generate_ppt_report(content: str, brand_name: str) -> str:
    """Generate a PowerPoint campaign deck using python-pptx (6 slides).

    Args:
        content: Markdown content to populate the presentation.
        brand_name: Brand name used in filename and title slide.

    Returns:
        Path to the saved .pptx file.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        ts   = datetime.now().strftime("%Y%m%d_%H%M%S")
        path = _REPORTS / f"campaign_report_{_safe_filename(brand_name)}_{ts}.md"
        path.write_text(content, encoding="utf-8")
        return f"[python-pptx not installed – saved as markdown: {path}]"

    ts        = datetime.now().strftime("%Y%m%d_%H%M%S")
    pptx_path = _REPORTS / f"campaign_report_{_safe_filename(brand_name)}_{ts}.pptx"

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    DARK_BG   = RGBColor(0x0D, 0x11, 0x17)
    ACCENT    = RGBColor(0x1A, 0x73, 0xE8)
    WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
    LIGHT_GREY= RGBColor(0xCC, 0xCC, 0xCC)

    blank_layout = prs.slide_layouts[6]  # blank

    def _add_slide(title_text: str, body_lines: list[str]) -> None:
        slide = prs.slides.add_slide(blank_layout)
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = DARK_BG

        # Title bar
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.9))
        tf = title_box.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = title_text
        p.alignment = PP_ALIGN.LEFT
        run = p.runs[0]
        run.font.bold  = True
        run.font.size  = Pt(28)
        run.font.color.rgb = ACCENT

        # Divider line (thin box)
        line = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            Inches(0.5), Inches(1.25), Inches(12.3), Inches(0.04),
        )
        line.fill.solid()
        line.fill.fore_color.rgb = ACCENT
        line.line.fill.background()

        # Body text
        body_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.7)
        )
        tf2 = body_box.text_frame
        tf2.word_wrap = True
        first = True
        for line_text in body_lines[:18]:  # cap at 18 lines per slide
            para = tf2.paragraphs[0] if first else tf2.add_paragraph()
            para.text = line_text
            if line_text.startswith("• "):
                para.level = 1
            run2 = para.runs[0] if para.runs else para.add_run()
            run2.font.size  = Pt(14)
            run2.font.color.rgb = LIGHT_GREY
            first = False

    # ── Extract sections from markdown ───────────────────────────────────────
    sections: dict[str, list[str]] = {}
    current_heading = "Overview"
    for raw_line in content.splitlines():
        line = raw_line.strip()
        if line.startswith("## "):
            current_heading = line[3:]
            sections.setdefault(current_heading, [])
        elif line.startswith("### "):
            sections.setdefault(current_heading, []).append(f"▸ {line[4:]}")
        elif line.startswith("- ") or line.startswith("* "):
            sections.setdefault(current_heading, []).append(f"• {line[2:]}")
        elif line and not line.startswith("#"):
            sections.setdefault(current_heading, []).append(line)

    # ── Slide 1: Title ────────────────────────────────────────────────────────
    slide1 = prs.slides.add_slide(blank_layout)
    s1bg = slide1.background.fill
    s1bg.solid()
    s1bg.fore_color.rgb = DARK_BG

    title_box = slide1.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"Digital Marketing Campaign"
    p.alignment = PP_ALIGN.CENTER
    p.runs[0].font.bold = True
    p.runs[0].font.size = Pt(36)
    p.runs[0].font.color.rgb = WHITE

    sub_box = slide1.shapes.add_textbox(Inches(1), Inches(3.7), Inches(11), Inches(1))
    tf2 = sub_box.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = brand_name
    p2.alignment = PP_ALIGN.CENTER
    p2.runs[0].font.size = Pt(24)
    p2.runs[0].font.color.rgb = ACCENT

    date_box = slide1.shapes.add_textbox(Inches(1), Inches(4.8), Inches(11), Inches(0.5))
    tf3 = date_box.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = datetime.now().strftime("%B %Y")
    p3.alignment = PP_ALIGN.CENTER
    p3.runs[0].font.size = Pt(16)
    p3.runs[0].font.color.rgb = LIGHT_GREY

    # ── Slides 2–6: section content ───────────────────────────────────────────
    desired_sections = [
        "Executive Summary", "Campaign Objectives", "Strategy Overview",
        "Content Plan Summary", "KPIs & Success Metrics",
    ]
    slide_count = 0
    for heading in desired_sections:
        # find closest match in sections dict
        matched_key = next(
            (k for k in sections if heading.lower() in k.lower()), None
        )
        body_lines = sections.get(matched_key, [f"See full report for details on {heading}."]) \
            if matched_key else [f"See full report for details on {heading}."]

        _add_slide(heading, body_lines)
        slide_count += 1
        if slide_count >= 5:
            break

    prs.save(str(pptx_path))
    print(f"[ReportTools] PPTX report saved → {pptx_path}", flush=True)
    return str(pptx_path)
